# ppt_operations.py - PPT generation helpers
from pathlib import Path
from typing import Any, Self

import pandas as pd
import pptx_ea_font
from loguru import logger
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_DATA_LABEL_POSITION, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.slide import Slide
from pptx.util import Cm, Pt

from config import CHART_THEMES
from utils import parse_markdown_style

from .layout_manager import layout_manager
from .schemas import (
    AxisChartConfig,
    BarChartConfig,
    BaseChartConfig,
    LayoutModel,
    LayoutType,
    LineChartConfig,
    PieChartConfig,
    RectangleStyleModel,
    TableConfig,
    TextContentModel,
)


class PPTOperations:
    """Low-level PPT operations used by the rendering pipeline."""

    DEFAULT_SLIDE_WIDTH_CM = 25.4
    DEFAULT_SLIDE_HEIGHT_CM = 14.29
    BLANK_LAYOUT_INDEX = 6

    def __init__(
        self, output_file_path: str | Path, template_file_path: str | Path | None = None
    ):
        """Create a PPT operations helper."""
        self.output_path = Path(output_file_path)
        self.template_path = Path(template_file_path) if template_file_path else None
        self.presentation = self._load_presentation()

    def _load_presentation(self) -> Presentation:
        """Load a template, an existing file, or a blank presentation."""
        if self.template_path and self.template_path.exists():
            logger.info(f"Loading template: {self.template_path}")
            return Presentation(str(self.template_path))

        if self.output_path.exists():
            logger.info(f"Loading existing PPT: {self.output_path}")
            return Presentation(str(self.output_path))

        logger.info("Creating new blank Presentation")
        return Presentation()

    def __enter__(self) -> Self:
        return self

    def __exit__(self, exc_type, exc_val, exc_tb):
        if exc_type is None:
            self.save()
        else:
            logger.error(f"PPT operations exited with error: {exc_val}")

    def save(self) -> None:
        """Save the presentation to disk."""
        self.output_path.parent.mkdir(parents=True, exist_ok=True)
        try:
            self.presentation.save(str(self.output_path))
            logger.info(f"PPT saved to: {self.output_path}")
        except Exception as e:
            logger.error(f"Failed to save PPT: {e}")
            raise

    def _get_slide(self, page_number: int) -> Slide:
        """Return the target slide using a 1-based page number."""
        total_slides = len(self.presentation.slides)
        if not (1 <= page_number <= total_slides):
            raise IndexError(
                f"Slide {page_number} is out of range (total: {total_slides})"
            )
        return self.presentation.slides[page_number - 1]

    def init_slides(
        self,
        slide_count: int,
        slide_width_cm: float = None,
        slide_height_cm: float = None,
        layout_type: LayoutType | str = None,
    ) -> None:
        """Initialize the deck size and create missing blank slides."""
        if layout_type and not slide_width_cm and not slide_height_cm:
            layout_str = (
                layout_type.value
                if isinstance(layout_type, LayoutType)
                else layout_type
            )
            slide_size = layout_manager.get_slide_size(layout_str)
            width_cm = slide_size.width
            height_cm = slide_size.height
        else:
            width_cm = slide_width_cm or self.DEFAULT_SLIDE_WIDTH_CM
            height_cm = slide_height_cm or self.DEFAULT_SLIDE_HEIGHT_CM

        self.presentation.slide_width = Cm(width_cm)
        self.presentation.slide_height = Cm(height_cm)

        current_slide_count = len(self.presentation.slides)
        slides_needed = slide_count - current_slide_count
        blank_layout = self.presentation.slide_layouts[self.BLANK_LAYOUT_INDEX]

        for _ in range(slides_needed):
            self.presentation.slides.add_slide(blank_layout)

        logger.info(
            f"PPT initialized with {len(self.presentation.slides)} slides "
            f"({width_cm}x{height_cm}cm)"
        )

    def add_text_box(
        self, page_num: int, content: TextContentModel, layout: LayoutModel
    ) -> None:
        """Add a formatted text box to a slide."""
        slide = self._get_slide(page_num)
        shape = slide.shapes.add_textbox(
            Cm(layout.left), Cm(layout.top), Cm(layout.width), Cm(layout.height or 1.0)
        )

        tf = shape.text_frame
        tf.word_wrap = content.word_wrap

        p = tf.paragraphs[0]
        p.alignment = layout.alignment.pptx_val
        p.clear()

        for segment in parse_markdown_style(content.text):
            run = p.add_run()
            run.text = segment.text
            run.font.size = Pt(content.font_size)
            run.font.color.rgb = content.font_color.rgb
            try:
                run.font.name = content.font_name
            except Exception as e:
                logger.warning(
                    f"Failed to use font {content.font_name}: {e}; fallback will apply."
                )
            run.font.bold = segment.is_bold or content.font_bold

        logger.info(f"Page {page_num}: added text box '{content.text[:50]}...'")

    def add_table(
        self,
        page_num: int,
        layout: LayoutModel,
        data: pd.DataFrame,
        config: TableConfig | None = None,
    ) -> None:
        """Add a table with adaptive row heights and column widths."""
        if config is None:
            config = TableConfig()

        data_reset = data
        rows, cols = data_reset.shape
        total_rows = rows + 1

        slide = self._get_slide(page_num)
        header_font_size, body_font_size = self._resolve_table_font_sizes(
            layout, config, rows
        )

        table_shape = slide.shapes.add_table(
            total_rows,
            cols,
            Cm(layout.left),
            Cm(layout.top),
            Cm(layout.width),
            Cm(layout.height),
        )
        table = table_shape.table
        self._apply_table_row_heights(table, layout.height, rows)
        self._apply_table_column_widths(table, data_reset, layout.width)

        def _set_cell_text(cell, text, font_size, font_color, bg_color, is_bold=False):
            cell.text = str(text)
            cell.margin_left = Cm(config.cell_margin_cm)
            cell.margin_right = Cm(config.cell_margin_cm)
            cell.margin_top = Cm(config.cell_margin_cm)
            cell.margin_bottom = Cm(config.cell_margin_cm)
            cell.text_frame.word_wrap = True

            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(font_size)
                    run.font.color.rgb = font_color.rgb
                    run.font.bold = is_bold
                    try:
                        pptx_ea_font.set_font(run, config.font_name)
                    except Exception:
                        logger.warning(
                            f"Failed to apply East Asian font {config.font_name}; "
                            "fallback font will be used."
                        )

            cell.fill.solid()
            cell.fill.fore_color.rgb = bg_color.rgb

        for col_idx, col_name in enumerate(data_reset.columns):
            _set_cell_text(
                table.cell(0, col_idx),
                str(col_name),
                header_font_size,
                config.header_font_color,
                config.header_bg_color,
                config.header_font_bold,
            )

        for row_idx in range(rows):
            for col_idx in range(cols):
                val = data_reset.iloc[row_idx, col_idx]
                _set_cell_text(
                    table.cell(row_idx + 1, col_idx),
                    str(val),
                    body_font_size,
                    config.body_font_color,
                    config.body_bg_color,
                    config.body_font_bold,
                )

        logger.debug(f"Page {page_num}: added {rows}x{cols} table")

    def _resolve_table_font_sizes(
        self, layout: LayoutModel, config: TableConfig, body_rows: int
    ) -> tuple[int, int]:
        """Shrink table fonts when the row count gets high."""
        header_row_factor = 1.15
        total_units = max(body_rows + header_row_factor, 1)
        body_row_height_cm = layout.height / total_units
        header_row_height_cm = body_row_height_cm * header_row_factor

        header_cap = max(
            int(header_row_height_cm * 28.346 * 0.5),
            config.min_header_font_size,
        )
        body_cap = max(
            int(body_row_height_cm * 28.346 * 0.5),
            config.min_body_font_size,
        )

        header_font_size = max(
            config.min_header_font_size,
            min(config.header_font_size, header_cap),
        )
        body_font_size = max(
            config.min_body_font_size,
            min(config.body_font_size, body_cap),
        )
        return header_font_size, body_font_size

    def _apply_table_row_heights(
        self, table, table_height_cm: float, body_rows: int
    ) -> None:
        """Keep every row inside the configured table frame."""
        if len(table.rows) == 0:
            return

        header_row_factor = 1.15
        total_units = max(body_rows + header_row_factor, 1)
        body_row_height_cm = table_height_cm / total_units
        header_row_height_cm = body_row_height_cm * header_row_factor

        table.rows[0].height = Cm(header_row_height_cm)
        for row_idx in range(1, len(table.rows)):
            table.rows[row_idx].height = Cm(body_row_height_cm)

    def _apply_table_column_widths(
        self, table, data: pd.DataFrame, total_width_cm: float
    ) -> None:
        """Spread table width by content length instead of a fixed first-column ratio."""
        cols = len(table.columns)
        if cols == 0:
            return
        if cols == 1:
            table.columns[0].width = Cm(total_width_cm)
            return

        min_col_width_cm = min(
            max(total_width_cm * 0.045, 0.75),
            total_width_cm / cols,
        )
        guaranteed_width = min_col_width_cm * cols
        if guaranteed_width >= total_width_cm:
            equal_width = total_width_cm / cols
            for col_idx in range(cols):
                table.columns[col_idx].width = Cm(equal_width)
            return

        weights: list[float] = []
        for col_idx, col_name in enumerate(data.columns):
            header_len = len(str(col_name).strip())
            series = data.iloc[:, col_idx].astype(str)
            max_value_len = (
                series.map(lambda value: len(value.strip())).max()
                if not series.empty
                else 0
            )
            weight = max(header_len * 1.15, max_value_len)
            if col_idx == 0:
                weight *= 1.2
            weights.append(max(weight, 1.0))

        extra_width = total_width_cm - guaranteed_width
        weight_sum = sum(weights) or cols

        for col_idx, weight in enumerate(weights):
            width_cm = min_col_width_cm + extra_width * (weight / weight_sum)
            table.columns[col_idx].width = Cm(width_cm)

    def add_rectangle(
        self,
        page_num: int,
        layout: LayoutModel,
        style: RectangleStyleModel | None = None,
    ) -> None:
        """Add a rectangle shape to a slide."""
        slide = self._get_slide(page_num)

        if style is None:
            style = RectangleStyleModel()

        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Cm(layout.left),
            Cm(layout.top),
            Cm(layout.width),
            Cm(layout.height),
        )
        shape.rotation = style.rotation

        shape.fill.solid()
        shape.fill.fore_color.rgb = style.fore_color.rgb

        if style.line_width > 0:
            shape.line.color.rgb = style.line_color.rgb
            shape.line.width = Pt(style.line_width)
        else:
            shape.line.fill.background()

        if style.is_background:
            shape.fill.background()

        logger.info(f"Page {page_num}: added rectangle")

    def add_picture(self, page_num: int, img_path: str, layout: LayoutModel) -> None:
        """Add an image to a slide."""
        slide = self._get_slide(page_num)
        path = Path(img_path)
        if not path.exists():
            logger.warning(f"Image not found: {path}")
            return

        slide.shapes.add_picture(
            str(path),
            Cm(layout.left),
            Cm(layout.top),
            width=Cm(layout.width),
        )
        logger.info(f"Page {page_num}: added picture {path.name}")

    def _prepare_chart_data(self, df: pd.DataFrame) -> ChartData:
        """Convert a DataFrame into python-pptx ChartData."""
        chart_data = ChartData()
        chart_data.categories = list(df.columns)

        for i in range(len(df)):
            series_data = df.iloc[i].fillna(0).tolist()
            series_name = str(df.index[i])
            chart_data.add_series(series_name, series_data)

        return chart_data

    def _apply_chart_formatting(self, chart, config: BaseChartConfig) -> None:
        """Apply shared chart formatting."""
        if hasattr(chart, "font"):
            chart.font.name = config.font_name
            chart.font.size = Pt(config.font_size)

        chart.has_legend = config.has_legend
        if config.has_legend:
            chart.legend.position = XL_LEGEND_POSITION.BOTTOM
            chart.legend.include_in_layout = False
            chart.legend.font.size = Pt(config.font_size)

        chart.has_title = bool(config.title)
        if chart.has_title:
            chart.chart_title.text_frame.text = config.title

        if config.style_name in CHART_THEMES:
            colors = CHART_THEMES[config.style_name]
            for i, series in enumerate(chart.series):
                if i < len(colors):
                    if hasattr(series.format.fill, "solid"):
                        series.format.fill.solid()
                        series.format.fill.fore_color.rgb = colors[i]
                    if hasattr(series.format.line, "color"):
                        series.format.line.color.rgb = colors[i]

    def _create_base_chart(
        self,
        page_num: int,
        layout: LayoutModel,
        df: pd.DataFrame,
        chart_type: XL_CHART_TYPE,
    ):
        """Create a chart and return the chart object."""
        slide = self._get_slide(page_num)
        chart_data = self._prepare_chart_data(df)

        graphic_frame = slide.shapes.add_chart(
            chart_type,
            Cm(layout.left),
            Cm(layout.top),
            Cm(layout.width),
            Cm(layout.height),
            chart_data,
        )
        return graphic_frame.chart

    def _configure_axes(self, chart: Any, config: AxisChartConfig) -> None:
        """Apply shared axis configuration for bar and line charts."""
        try:
            val_axis = chart.value_axis
            val_axis.visible = config.y_axis_visible
            val_axis.has_major_gridlines = False
            val_axis.tick_labels.font.size = Pt(config.font_size)
            if config.value_axis_max:
                val_axis.maximum_scale = config.value_axis_max
            if config.value_axis_format:
                val_axis.tick_labels.number_format = config.value_axis_format

            cat_axis = chart.category_axis
            cat_axis.visible = config.x_axis_visible
            cat_axis.tick_labels.font.size = Pt(config.font_size)
        except ValueError:
            pass

    def add_bar_chart(
        self,
        page_num: int,
        df: pd.DataFrame,
        layout: LayoutModel,
        config: BarChartConfig,
    ) -> None:
        """Add a bar chart."""
        c_type = (
            XL_CHART_TYPE.COLUMN_STACKED
            if config.grouping == "stacked"
            else XL_CHART_TYPE.COLUMN_CLUSTERED
        )

        chart = self._create_base_chart(page_num, layout, df, c_type)
        self._apply_chart_formatting(chart, config)
        self._configure_axes(chart, config)

        if chart.plots:
            plot = chart.plots[0]
            plot.gap_width = config.gap_width
            plot.overlap = config.overlap
            plot.vary_by_categories = False

            if config.has_data_labels:
                plot.has_data_labels = True
                labels = plot.data_labels
                labels.font.size = Pt(config.font_size - 1)
                labels.position = XL_DATA_LABEL_POSITION.OUTSIDE_END
            else:
                plot.has_data_labels = False

        logger.debug(f"Page {page_num}: Added Bar Chart")

    def add_line_chart(
        self,
        page_num: int,
        df: pd.DataFrame,
        layout: LayoutModel,
        config: LineChartConfig,
    ) -> None:
        """Add a line chart."""
        c_type = (
            XL_CHART_TYPE.LINE_MARKERS if config.has_markers else XL_CHART_TYPE.LINE
        )

        chart = self._create_base_chart(page_num, layout, df, c_type)
        self._apply_chart_formatting(chart, config)
        self._configure_axes(chart, config)

        if chart.plots:
            plot = chart.plots[0]
            if config.has_data_labels:
                plot.has_data_labels = True
                data_labels = plot.data_labels
                data_labels.font.size = Pt(config.font_size)
                data_labels.position = XL_DATA_LABEL_POSITION.ABOVE
            else:
                plot.has_data_labels = False

            for series in chart.series:
                series.smooth = config.smooth_line
                series.format.line.width = Pt(config.line_width)

        logger.debug(f"Page {page_num}: Added Line Chart")

    def add_pie_chart(
        self,
        page_num: int,
        df: pd.DataFrame,
        layout: LayoutModel,
        config: PieChartConfig,
    ) -> None:
        """Add a pie or doughnut chart."""
        c_type = XL_CHART_TYPE.DOUGHNUT if config.hole_size > 0 else XL_CHART_TYPE.PIE
        chart = self._create_base_chart(page_num, layout, df, c_type)
        self._apply_chart_formatting(chart, config)

        if chart.plots:
            plot = chart.plots[0]
            plot.vary_by_categories = True
            if config.has_data_labels:
                plot.has_data_labels = True
                labels = plot.data_labels
                labels.font.size = Pt(max(config.font_size - 1, 8))
                labels.show_category_name = config.show_category_name
                labels.show_percentage = config.show_percentage
                labels.show_value = config.show_value
                labels.number_format = "0%" if config.show_percentage else "General"
            else:
                plot.has_data_labels = False

        if config.style_name in CHART_THEMES and chart.series:
            colors = CHART_THEMES[config.style_name]
            for idx, point in enumerate(chart.series[0].points):
                point.format.fill.solid()
                point.format.fill.fore_color.rgb = colors[idx % len(colors)]

        logger.debug(f"Page {page_num}: Added Pie Chart")
