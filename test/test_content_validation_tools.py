from __future__ import annotations

import tempfile
import unittest
from pathlib import Path
from typing import Any

import pandas as pd
import yaml
from langchain.tools import ToolRuntime
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches

from method.agents.content_validation.agent import build_content_payload
from method.agents.content_validation.artifacts import write_content_artifacts
from method.agents.content_validation.tools import (
    CONTENT_VALIDATION_TOOLS,
    align_visible_dataframe,
    compare_display_dataframes,
    execute_table_state,
)
from method.schemas import TableAnalysisConfig
from method.transformers import StatTransformer


class FakeDatabase:
    def query(self, sql, params=None):
        del sql, params
        return pd.DataFrame(
            {
                "date_code": ["2020-01-01", "2021-01-01", "2021-02-01"],
                "trade_sets": [1, 1, 1],
            }
        )


class FakeClient:
    def respond(self, request):
        del request
        return {"response": "Yes, please apply the proposed update."}


class FakeToolContext:
    client = FakeClient()


class ContentValidationToolsTest(unittest.TestCase):
    def test_table_analysis_config_rejects_misspelled_table_type(self) -> None:
        with self.assertRaises(ValueError):
            TableAnalysisConfig.model_validate(
                {
                    "table_type": "constraint-filed",
                    "dimensions": [],
                    "metrics": [],
                }
            )

    def test_build_content_payload_includes_editable_element_ids(self) -> None:
        payload = build_content_payload(_analysis_state(data_path="visible.csv"))

        self.assertEqual(payload["title"]["element_id"], "1")
        self.assertEqual(payload["summary"]["element_id"], "2")
        self.assertEqual(payload["tables"][0]["caption"]["element_id"], "3")
        self.assertEqual(payload["tables"][0]["body"]["element_id"], "4")
        self.assertNotIn("scope_dialogue", payload)

    def test_compare_display_dataframes_allows_rounding_noise(self) -> None:
        visible = pd.DataFrame({"year": [2020], "value": [100.0]})
        expected = pd.DataFrame({"year": [2020], "value": [100.4]})

        result = compare_display_dataframes(visible, expected)

        self.assertEqual(result["status"], "equal")

    def test_compare_display_dataframes_reports_value_differences(self) -> None:
        visible = pd.DataFrame(
            {"year": [2020, 2021], "trade_counts": [10, 20], "amount": [1.0, 2.0]}
        )
        expected = pd.DataFrame(
            {"year": [2020, 2021], "trade_counts": [11, 25], "amount": [1.0, 9.0]}
        )

        result = compare_display_dataframes(visible, expected)

        self.assertEqual(result["status"], "different")
        self.assertEqual(result["diff_count"], 3)
        self.assertEqual(
            result["diff_examples"],
            [
                {"row": 0, "column": "trade_counts", "visible": 10, "expected": 11},
                {"row": 1, "column": "trade_counts", "visible": 20, "expected": 25},
                {"row": 1, "column": "amount", "visible": 2.0, "expected": 9.0},
            ],
        )
        self.assertIn("3 cell(s) differ", result["diff_summary"])

    def test_compare_display_dataframes_ignores_top_left_header_cell(self) -> None:
        visible = pd.DataFrame(
            {
                "__index__": ["Supply Count", "Sales Count"],
                "0-20m²": [1004, 786],
                "20-40m²": [212, 522],
            }
        )
        expected = pd.DataFrame(
            {
                "metric": ["Supply Count", "Sales Count"],
                "0-20m²": [1004, 786],
                "20-40m²": [212, 522],
            }
        )

        result = compare_display_dataframes(visible, expected)

        self.assertEqual(result["status"], "equal")

    def test_align_visible_dataframe_maps_chart_category_to_dimension(self) -> None:
        visible = pd.DataFrame(
            {
                "category": ["2020-01", "2020-02"],
                "trade_counts": [11.0, 7.0],
            }
        )
        expected = pd.DataFrame(
            {
                "month": ["2020-01", "2020-02"],
                "trade_counts": [11, 7],
            }
        )
        state = _analysis_state(data_path="visible.csv")
        state["tables"][0]["body"]["type"] = "chart-line"
        state["tables"][0]["calculation_logic"]["dimensions"][0]["target_col"] = "month"

        aligned = align_visible_dataframe(visible, state["tables"][0])
        result = compare_display_dataframes(aligned, expected)

        self.assertEqual(list(aligned.columns), ["month", "trade_counts"])
        self.assertEqual(result["status"], "equal")

    def test_modify_tools_return_agent_state_updates(self) -> None:
        chart_state = _content_validation_state(
            _analysis_state(data_path="visible.csv")
        )
        chart_update = _invoke_content_tool(
            "modify_chart",
            chart_state,
            {"table_index": 0, "data_path": "computed.csv"},
        ).update

        self.assertEqual(
            chart_update["analysis_state"]["tables"][0]["body"]["data_path"],
            "computed.csv",
        )
        self.assertEqual(chart_update["tool_log"][0]["tool"], "modify_chart")
        self.assertEqual(chart_update["tool_log"][0]["result"], {"success": True})

        table_analysis_state = _analysis_state(data_path="visible.csv")
        table_analysis_state["tables"][0]["body"]["type"] = "table"
        table_state = _content_validation_state(table_analysis_state)
        table_update = _invoke_content_tool(
            "modify_table",
            table_state,
            {"table_index": 0, "data_path": "computed.csv"},
        ).update

        self.assertEqual(
            table_update["analysis_state"]["tables"][0]["body"]["data_path"],
            "computed.csv",
        )
        self.assertEqual(table_update["tool_log"][0]["tool"], "modify_table")
        self.assertEqual(table_update["tool_log"][0]["result"], {"success": True})

        textbox_state = _content_validation_state(
            _analysis_state(data_path="visible.csv")
        )
        textbox_update = _invoke_content_tool(
            "modify_textbox",
            textbox_state,
            {"element_id": "3", "text": "Updated caption"},
        ).update

        self.assertEqual(
            textbox_update["analysis_state"]["tables"][0]["caption"]["text"],
            "Updated caption",
        )
        self.assertEqual(textbox_update["tool_log"][0]["tool"], "modify_textbox")
        self.assertEqual(textbox_update["tool_log"][0]["result"], {"success": True})

    def test_ask_client_records_detected_issue_labels(self) -> None:
        state = _content_validation_state(_analysis_state(data_path="visible.csv"))

        update = _invoke_content_tool(
            "ask_client",
            state,
            {
                "error_type": "claim_error",
                "description": "Caption says table but body is a bar chart.",
                "target": "st.caption",
            },
            context=FakeToolContext(),
        ).update

        self.assertEqual(
            update["detected_issues"],
            [
                {
                    "request_type": "content_update_confirmation",
                    "target": "st.caption",
                    "error_type": "claim_error",
                    "evidence": "Caption says table but body is a bar chart.",
                    "client_response": "Yes, please apply the proposed update.",
                    "confirmed": False,
                }
            ],
        )

    def test_execute_table_state_runs_extracted_state(self) -> None:
        expected = execute_table_state(
            _analysis_state(data_path="visible.csv")["tables"][0],
            transformer=StatTransformer(),
            query_func=FakeDatabase().query,
        )

        self.assertEqual(list(expected.columns), ["year", "trade_counts"])
        self.assertEqual(expected["trade_counts"].tolist(), [1, 2])

    def test_execute_table_state_adds_filter_columns_for_execution(self) -> None:
        captured_sql = []

        class PriceDatabase:
            def query(self, sql, params=None):
                del params
                captured_sql.append(sql)
                return pd.DataFrame(
                    {
                        "date_code": ["2020-01-01", "2020-02-01", "2021-01-01"],
                        "dim_unit_price": [100.0, 300.0, 500.0],
                        "trade_sets": [1, 0, 1],
                    }
                )

        state = _analysis_state(data_path="visible.csv")
        table_state = state["tables"][0]
        table_state["caption"]["data_source"]["select_columns"] = [
            "date_code",
            "dim_unit_price",
        ]
        table_state["calculation_logic"] = {
            "table_type": "field-constraint",
            "dimensions": [
                {
                    "source_col": "date_code",
                    "target_col": "year",
                    "method": "period",
                    "time_granularity": "year",
                }
            ],
            "metrics": [
                {
                    "name": "avg_unit_price",
                    "source_col": "dim_unit_price",
                    "agg_func": "mean",
                    "filter_condition": {"trade_sets": 1},
                }
            ],
        }

        expected = execute_table_state(
            table_state,
            transformer=StatTransformer(),
            query_func=PriceDatabase().query,
        )

        self.assertIn("trade_sets", captured_sql[0])
        self.assertEqual(list(expected.columns), ["year", "avg_unit_price"])
        self.assertEqual(expected["avg_unit_price"].tolist(), [100.0, 500.0])

    def test_execute_table_state_does_not_add_metric_columns_for_execution(self) -> None:
        captured_sql = []

        class IncompleteDatabase:
            def query(self, sql, params=None):
                del params
                captured_sql.append(sql)
                return pd.DataFrame(
                    {
                        "date_code": ["2020-01-01"],
                        "trade_sets": [1],
                    }
                )

        state = _analysis_state(data_path="visible.csv")
        table_state = state["tables"][0]
        table_state["caption"]["data_source"]["select_columns"] = [
            "date_code",
        ]
        table_state["calculation_logic"] = {
            "table_type": "field-constraint",
            "dimensions": [
                {
                    "source_col": "date_code",
                    "target_col": "year",
                    "method": "period",
                    "time_granularity": "year",
                }
            ],
            "metrics": [
                {
                    "name": "avg_unit_price",
                    "source_col": "dim_unit_price",
                    "agg_func": "mean",
                    "filter_condition": {"trade_sets": 1},
                }
            ],
        }

        with self.assertRaisesRegex(KeyError, "Missing metric source column"):
            execute_table_state(
                table_state,
                transformer=StatTransformer(),
                query_func=IncompleteDatabase().query,
            )

        self.assertIn("trade_sets", captured_sql[0])
        self.assertNotIn("dim_unit_price", captured_sql[0])

    def test_execute_table_state_runs_range_dimension_without_min_max(self) -> None:
        class RangeDatabase:
            def query(self, sql, params=None):
                del sql, params
                return pd.DataFrame(
                    {
                        "dim_area": [5, 18, 22, 43],
                        "trade_sets": [1, 1, 1, 1],
                    }
                )

        state = _analysis_state(data_path="visible.csv")
        state["tables"][0]["calculation_logic"] = {
            "table_type": "field-constraint",
            "dimensions": [
                {
                    "source_col": "dim_area",
                    "target_col": "area_range",
                    "method": "range",
                    "step": 20,
                    "format_str": "{}-{}m²",
                }
            ],
            "metrics": [
                {
                    "name": "Area Rng Stats",
                    "source_col": "trade_sets",
                    "agg_func": "count",
                    "filter_condition": {"trade_sets": 1},
                }
            ],
        }
        state["tables"][0]["caption"]["data_source"]["select_columns"] = [
            "dim_area",
            "trade_sets",
        ]

        expected = execute_table_state(
            state["tables"][0],
            transformer=StatTransformer(),
            query_func=RangeDatabase().query,
        )

        self.assertEqual(list(expected.columns), ["area_range", "Area Rng Stats"])
        self.assertEqual(
            expected["area_range"].tolist(), ["0-20m²", "20-40m²", "40-60m²"]
        )
        self.assertEqual(expected["Area Rng Stats"].tolist(), [2, 1, 1])

    def test_write_content_artifacts_emits_yaml_and_repaired_csv(self) -> None:
        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            expected_path = root / "expected.csv"
            visible_path = root / "visible.csv"
            pd.DataFrame({"year": [2020], "trade_counts": [2]}).to_csv(
                expected_path,
                index=False,
            )
            pd.DataFrame({"year": [2020], "trade_counts": [1]}).to_csv(
                visible_path,
                index=False,
            )
            source_pptx = root / "slide.pptx"
            presentation = Presentation()
            slide = presentation.slides.add_slide(presentation.slide_layouts[6])
            for text in ("Old title", "Old summary", "Old caption"):
                slide.shapes.add_textbox(
                    Inches(1), Inches(1), Inches(3), Inches(0.5)
                ).text = text
            chart_data = ChartData()
            chart_data.categories = [2020]
            chart_data.add_series("trade_counts", [1])
            slide.shapes.add_chart(
                XL_CHART_TYPE.COLUMN_CLUSTERED,
                Inches(1),
                Inches(2),
                Inches(5),
                Inches(3),
                chart_data,
            )
            presentation.save(source_pptx)

            artifacts = write_content_artifacts(
                source_pptx=source_pptx,
                analysis_state=_analysis_state(data_path=str(expected_path)),
                artifact_dir=root / "review_artifacts",
            )

            repaired_yaml = yaml.safe_load(
                Path(artifacts["yaml_path"]).read_text(encoding="utf-8")
            )
            repaired_data_path = Path(repaired_yaml["tables"][0]["body"]["data_path"])
            self.assertTrue(repaired_data_path.exists())
            self.assertEqual(
                pd.read_csv(repaired_data_path)["trade_counts"].tolist(),
                [2],
            )
            repaired_pptx = Presentation(artifacts["pptx_path"])
            repaired_slide = repaired_pptx.slides[0]
            self.assertEqual(repaired_slide.shapes[0].text, "Example")
            self.assertEqual(repaired_slide.shapes[1].text, "Example summary")
            self.assertEqual(repaired_slide.shapes[2].text, "Example caption")
            self.assertEqual(
                list(repaired_slide.shapes[3].chart.series[0].values), [2.0]
            )

    def test_write_content_artifacts_updates_ppt_table(self) -> None:
        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            data_path = root / "expected.csv"
            pd.DataFrame({"year": [2020, 2021], "trade_counts": [2, 3]}).to_csv(
                data_path, index=False
            )

            source_pptx = root / "slide.pptx"
            presentation = Presentation()
            slide = presentation.slides.add_slide(presentation.slide_layouts[6])
            for text in ("Old title", "Old summary", "Old caption"):
                slide.shapes.add_textbox(
                    Inches(1), Inches(1), Inches(3), Inches(0.5)
                ).text = text
            slide.shapes.add_table(
                3,
                2,
                Inches(1),
                Inches(2),
                Inches(5),
                Inches(3),
            )
            presentation.save(source_pptx)

            state = _analysis_state(data_path=str(data_path))
            state["tables"][0]["body"]["type"] = "table"
            artifacts = write_content_artifacts(
                source_pptx=source_pptx,
                analysis_state=state,
                artifact_dir=root / "review_artifacts",
            )

            repaired_table = (
                Presentation(artifacts["pptx_path"]).slides[0].shapes[3].table
            )
            values = [
                [repaired_table.cell(row, column).text for column in range(2)]
                for row in range(3)
            ]
            self.assertEqual(
                values,
                [["year", "trade_counts"], ["2020", "2"], ["2021", "3"]],
            )


def _analysis_state(data_path: str) -> dict[str, Any]:
    return {
        "title": {"text": "Example", "element_id": "1"},
        "summary": {
            "text": "Example summary",
            "element_id": "2",
            "data_source": {
                "connection": {"table": "beijing_new_house"},
                "filters": {
                    "city": "Beijing",
                    "block": "Liangxiang",
                    "start_date": "2020-01-01",
                    "end_date": "2021-12-31",
                },
            },
        },
        "final_data_source": {
            "connection": {"table": "beijing_new_house"},
            "filters": {
                "city": "Beijing",
                "block": "Liangxiang",
                "start_date": "2020-01-01",
                "end_date": "2021-12-31",
            },
        },
        "tables": [
            {
                "caption": {
                    "text": "Example caption",
                    "element_id": "3",
                    "data_source": {
                        "connection": {"table": "beijing_new_house"},
                        "select_columns": ["date_code", "trade_sets"],
                        "filters": {
                            "city": "Beijing",
                            "block": "Liangxiang",
                            "start_date": "2020-01-01",
                            "end_date": "2021-12-31",
                        },
                    },
                },
                "body": {
                    "element_id": "4",
                    "type": "chart-bar",
                    "data_path": data_path,
                },
                "data_path": data_path,
                "calculation_logic": {
                    "table_type": "field-constraint",
                    "dimensions": [
                        {
                            "source_col": "date_code",
                            "target_col": "year",
                            "method": "period",
                            "time_granularity": "year",
                        }
                    ],
                    "metrics": [
                        {
                            "name": "trade_counts",
                            "source_col": "trade_sets",
                            "agg_func": "count",
                            "filter_condition": {"trade_sets": 1},
                        }
                    ],
                    "crosstab_row": None,
                    "crosstab_col": None,
                },
            }
        ],
    }


def _content_validation_state(analysis_state: dict[str, Any]) -> dict[str, Any]:
    return {
        "messages": [],
        "analysis_state": analysis_state,
        "table_records": [],
        "tool_log": [],
        "detected_issues": [],
    }


def _invoke_content_tool(
    name: str,
    state: dict[str, Any],
    args: dict[str, Any],
    context: Any | None = None,
):
    if context is None:
        context = object()
    tools = {
        content_tool.name: content_tool for content_tool in CONTENT_VALIDATION_TOOLS
    }
    return tools[name].func(
        **args,
        runtime=ToolRuntime(
            state=state,
            context=context,
            config={},
            stream_writer=lambda _: None,
            tool_call_id=None,
            store=None,
        ),
    )


def _ppt_representation(data_path: str) -> dict[str, Any]:
    return {
        "title": {"text": "Example"},
        "summary": {"text": "Example summary"},
        "structured_tables": [
            {
                "caption": {"text": "Example caption"},
                "body": {"type": "chart-bar", "data_path": data_path},
            }
        ],
    }


if __name__ == "__main__":
    unittest.main()
