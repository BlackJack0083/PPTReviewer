from __future__ import annotations

import csv
import json
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_SHAPE_TYPE

from method.utils import Client, parse_json_object

from ..types import SlideReviewInput

PROMPT_DIR = Path(__file__).resolve().parents[2] / "prompts"
DEFAULT_ROLE_LABELING_PROMPT_PATH = PROMPT_DIR / "role_labeling_prompt.txt"

ROLE_SET = {
    "title",
    "summary",
    "caption",
    "chart-bar",
    "chart-line",
    "chart-pie",
    "table",
}


def get_shape(shape) -> str:
    """判断 PPTX shape 的粗粒度类型。

    Args:
        shape: python-pptx 的 shape 对象。

    Returns:
        shape 类型字符串。可能值包括 `text`、`table`、`chart-bar`、
        `chart-line`、`chart-pie`、`other`。
    """
    if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
        return "table"

    if shape.shape_type == MSO_SHAPE_TYPE.CHART and hasattr(shape, "chart"):
        chart_type = shape.chart.chart_type
        if chart_type in _chart_type_values(
            "COLUMN_CLUSTERED",
            "COLUMN_STACKED",
            "COLUMN_STACKED_100",
            "BAR_CLUSTERED",
            "BAR_STACKED",
            "BAR_STACKED_100",
        ):
            return "chart-bar"
        if chart_type in _chart_type_values(
            "LINE",
            "LINE_MARKERS",
            "LINE_STACKED",
            "LINE_STACKED_100",
            "LINE_MARKERS_STACKED",
            "LINE_MARKERS_STACKED_100",
        ):
            return "chart-line"
        if chart_type in _chart_type_values("PIE", "PIE_EXPLODED", "DOUGHNUT"):
            return "chart-pie"
        return "other"

    if shape.has_text_frame and shape.text.strip():
        return "text"

    return "other"


def extract_pptx_elements(pptx_path: Path, slide_idx: int = 0) -> dict[str, Any]:
    """从 PPTX 中抽取可编辑元素，不做语义判断。

    这一步只读取 PPTX 自身的结构信息，包括文本框、表格、图表和布局。
    不读取 benchmark gold 信息，也不抽取 chart/table 的二维数据。

    Args:
        pptx_path: 待解析的 PPTX 文件路径。
        slide_idx: 要解析的 slide 下标，当前流程默认只处理第 0 页。

    Returns:
        `observed_slide` 初始结构，包含 `slide_size` 和 `elements`。
        此时 `elements` 中还没有 `role` 字段。

    Raises:
        IndexError: 当 `slide_idx` 超出 PPTX 页数时抛出。
    """
    presentation = Presentation(pptx_path)
    if slide_idx >= len(presentation.slides):
        raise IndexError(f"Slide index {slide_idx} out of range.")

    def cm(emu: float) -> float:
        return round(emu / 360000.0, 2)

    slide = presentation.slides[slide_idx]
    elements: list[dict[str, Any]] = []

    for shape_idx, shape in enumerate(slide.shapes, 1):
        shape_type = get_shape(shape)
        if shape_type == "other":
            continue

        element_id = str(shape_idx)
        element: dict[str, Any] = {
            "id": element_id,
            "layout": {
                "x": cm(shape.left),
                "y": cm(shape.top),
                "width": cm(shape.width),
                "height": cm(shape.height),
            },
        }

        if shape_type == "text":
            element.update(
                {
                    "type": "textBox",
                    "text": shape.text.strip(),
                }
            )
        elif shape_type == "table":
            element["type"] = "table"
        else:
            element["type"] = "chart"
            element["shape"] = shape_type

        elements.append(element)

    return {
        "slide_size": {
            "width": cm(presentation.slide_width),
            "height": cm(presentation.slide_height),
        },
        "elements": elements,
    }


def validate_role_labels(
    observed_slide: dict[str, Any],
    roles: list[Any],
) -> list[dict[str, str]]:
    """校验 VLM role 标注是否完整且合法。

    Args:
        observed_slide: 无 role 的 PPTX 元素结构。
        roles: VLM 返回的 role 标注列表。

    Returns:
        规范化后的 role 标注列表，每项为 `{"id": ..., "role": ...}`。

    Raises:
        ValueError: role 标注未覆盖全部元素、包含重复元素或使用非法 role 时抛出。
    """
    expected_ids = {str(element["id"]) for element in observed_slide["elements"]}
    assignments = [
        {"id": str(item["id"]), "role": str(item["role"])}
        for item in roles
    ]
    assigned_ids = [item["id"] for item in assignments]
    if len(assigned_ids) != len(set(assigned_ids)) or set(assigned_ids) != expected_ids:
        raise ValueError(
            f"Role assignment ids must match slide element ids: "
            f"expected={sorted(expected_ids)}, actual={sorted(assigned_ids)}"
        )
    invalid_roles = sorted({item["role"] for item in assignments} - ROLE_SET)
    if invalid_roles:
        raise ValueError(f"Unsupported roles: {invalid_roles}")
    return assignments


class SlideParserAgent:
    """将 PPTX/PNG 解析为后续 verification 可读的 slide 表示。

    Parser 负责三件事：抽取 PPTX 元素、调用 VLM 标注 role、导出
    chart/table 数据 CSV，并组装 `ppt_representation`。
    """

    def __init__(
        self,
        *,
        model: str | None = None,
        api_key: str | None = None,
        base_url: str = "https://dashscope.aliyuncs.com/compatible-mode/v1",
        timeout_sec: int = 120,
        enable_thinking: bool | None = False,
        client: Any | None = None,
    ):
        """初始化 slide parser。

        Args:
            model: OpenAI-compatible VLM 模型名。未传入 `client` 时必须提供。
            api_key: API key；为空时由底层 `Client` 从环境变量读取。
            base_url: OpenAI-compatible API base URL。
            timeout_sec: 单次请求超时时间，单位秒。
            enable_thinking: 是否通过 `extra_body` 打开模型 thinking。
            client: 可选的已构造 client，主要用于测试或外部统一注入。

        Returns:
            None。

        Raises:
            ValueError: 当未提供 `client` 且 `model` 为空时抛出。
        """
        self.role_labeling_prompt = DEFAULT_ROLE_LABELING_PROMPT_PATH.read_text(
            encoding="utf-8"
        )
        if client is not None:
            self.client = client
        else:
            if model is None:
                raise ValueError("SlideParserAgent requires model when client is not provided.")
            self.client = Client(
                model=model,
                api_key=api_key,
                base_url=base_url,
                timeout_sec=timeout_sec,
                enable_thinking=enable_thinking,
            )

    async def arun(self, slide_input: SlideReviewInput) -> dict[str, Any]:
        """执行单页 slide 的 Phase 1 parser 流程。

        Args:
            slide_input: 当前 case 的输入路径，只包含 `pptx_path` 和
                `image_path`。

        Returns:
            包含两个顶层字段的字典：
            `observed_slide` 为带 role 的可编辑元素列表；
            `ppt_representation` 为简化后的 slide 表示，并引用导出的 CSV。
        """
        raw_elements = extract_pptx_elements(slide_input.pptx_path)
        prompt = "Input elements:\n" + json.dumps(
            {
                "slide_size": raw_elements["slide_size"],
                "elements": [
                    {
                        "id": element["id"],
                        "type": element["type"],
                        "layout": element["layout"],
                        **(
                            {"text": element["text"]}
                            if element["type"] == "textBox"
                            else {"shape": element.get("shape", element["type"])}
                        ),
                    }
                    for element in raw_elements["elements"]
                ],
                "allowed_roles": sorted(ROLE_SET),
            },
            ensure_ascii=False,
            indent=2,
        )
        content = await self.client.achat(
            self.role_labeling_prompt,
            prompt,
            image_path=slide_input.image_path,
            response_format="json_object",
        )
        roles = parse_json_object(content).get("roles")
        if not isinstance(roles, list):
            raise ValueError(f"VLM role response must contain roles list: {content}")

        role_by_id = {
            item["id"]: item["role"]
            for item in validate_role_labels(raw_elements, roles)
        }
        observed_slide = {
            "slide_size": raw_elements["slide_size"],
            "elements": [
                {
                    **{
                        key: value
                        for key, value in element.items()
                        if not str(key).startswith("_")
                    },
                    "role": role_by_id[element["id"]],
                }
                for element in raw_elements["elements"]
            ],
        }
        ppt_representation = build_ppt_representation(
            pptx_path=slide_input.pptx_path,
            observed_slide=observed_slide,
        )
        return {
            "observed_slide": observed_slide,
            "ppt_representation": ppt_representation,
        }


def build_ppt_representation(
    *,
    pptx_path: Path,
    observed_slide: dict[str, Any],
    slide_idx: int = 0,
) -> dict[str, Any]:
    """把带 role 的元素整理为紧凑的 PPT 表示。

    这一步只做结构整理和 chart/table 数据导出，不做业务分析、不推断
    数据库字段、不判断 slide 是否有错误。

    Args:
        pptx_path: PPTX 文件路径，CSV 会写入该文件所在目录。
        observed_slide: 已合并 role 的 slide 元素结构。
        slide_idx: 要读取 chart/table 数据的 slide 下标。

    Returns:
        `ppt_representation` 字典，包含 `title`、`summary` 和
        `structured_tables`。每个 table body 通过 `data_path` 指向 CSV。

    Raises:
        ValueError: 未找到 chart/table body 时抛出。
    """
    presentation = Presentation(pptx_path)
    slide = presentation.slides[slide_idx]
    elements = observed_slide["elements"]
    title = next(element for element in elements if element["role"] == "title")
    summary = next(element for element in elements if element["role"] == "summary")
    captions = [element for element in elements if element["role"] == "caption"]
    bodies = [
        element
        for element in elements
        if element["role"] in {"chart-bar", "chart-line", "chart-pie", "table"}
    ]

    if not bodies:
        raise ValueError("Parser found no chart/table body elements.")

    structured_tables = []
    output_dir = pptx_path.parent
    for table_idx, body_element in enumerate(bodies, 1):
        shape = slide.shapes[int(body_element["id"]) - 1]
        caption = _nearest_caption(body_element, captions)
        header, rows = _extract_body_table(shape=shape, role=body_element["role"])
        csv_path = output_dir / ("data.csv" if len(bodies) == 1 else f"data_{table_idx}.csv")
        with csv_path.open("w", encoding="utf-8", newline="") as handle:
            writer = csv.writer(handle)
            writer.writerow(header)
            writer.writerows(rows)
        structured_tables.append(
            {
                "caption": {
                    "element_id": str(caption["id"]),
                    "text": caption["text"],
                },
                "header": header,
                "body": {
                    "element_id": str(body_element["id"]),
                    "type": body_element["role"],
                    "data_path": str(csv_path),
                },
            }
        )

    return {
        "title": {"element_id": str(title["id"]), "text": title["text"]},
        "summary": {"element_id": str(summary["id"]), "text": summary["text"]},
        "structured_tables": structured_tables,
    }


def _nearest_caption(
    body_element: dict[str, Any],
    captions: list[dict[str, Any]],
) -> dict[str, Any]:
    """为 chart/table body 选择空间位置最近的 caption。

    Args:
        body_element: role 为 chart/table 的 body 元素。
        captions: 候选 caption 元素列表。

    Returns:
        距离 body 中心点最近的 caption。

    Raises:
        ValueError: 没有 caption 候选时抛出。
    """
    if not captions:
        raise ValueError(f"Parser found no caption for body element {body_element['id']}.")

    def center(layout: dict[str, Any]) -> tuple[float, float]:
        return (
            float(layout.get("x", 0.0)) + float(layout.get("width", 0.0)) / 2.0,
            float(layout.get("y", 0.0)) + float(layout.get("height", 0.0)) / 2.0,
        )

    body_center = center(body_element.get("layout", {}))
    return min(
        captions,
        key=lambda caption: abs(body_center[0] - center(caption.get("layout", {}))[0])
        + abs(body_center[1] - center(caption.get("layout", {}))[1]),
    )


def _extract_body_table(shape, role: str) -> tuple[list[str], list[list[Any]]]:
    """从 table 或 chart shape 中抽取二维数据。

    Args:
        shape: python-pptx 的 table/chart shape 对象。
        role: body 元素 role，用于区分 table 与 chart。

    Returns:
        二元组 `(header, rows)`。`header` 是 CSV 表头，`rows` 是数据行。
    """
    if role == "table":
        raw_rows = [
            [
                str(shape.table.cell(row_idx, col_idx).text).strip()
                for col_idx in range(len(shape.table.columns))
            ]
            for row_idx in range(len(shape.table.rows))
        ]
        if not raw_rows or not raw_rows[0]:
            raise ValueError("Parser cannot export an empty PPT table.")
        return raw_rows[0], raw_rows[1:]

    chart = shape.chart
    categories = [str(category.label) for category in chart.plots[0].categories]
    series = list(chart.series)
    if not categories or not series:
        raise ValueError("Parser cannot export an empty chart.")
    header = ["category"] + [str(item.name) for item in series]
    rows = []
    for idx, category in enumerate(categories):
        rows.append(
            [category] + [item.values[idx] if idx < len(item.values) else "" for item in series]
        )
    return header, rows


def _chart_type_values(*names: str) -> set[Any]:
    """把 chart type 常量名转换为 python-pptx 枚举值集合。

    Args:
        *names: `XL_CHART_TYPE` 上的常量名。

    Returns:
        存在于当前 python-pptx 版本中的枚举值集合。
    """
    values = set()
    for name in names:
        if hasattr(XL_CHART_TYPE, name):
            values.add(getattr(XL_CHART_TYPE, name))
    return values
