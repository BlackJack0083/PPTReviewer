from __future__ import annotations

import csv
import json
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_SHAPE_TYPE

from method.utils import Client, parse_json_object

from .types import SlideReviewInput

PROMPT_DIR = Path(__file__).resolve().parents[1] / "prompts"
DEFAULT_ROLE_LABELING_PROMPT_PATH = PROMPT_DIR / "role_labeling_prompt.txt"

ROLE_SET = {
    "title",
    "summary",
    "caption",
    "chart-bar",
    "chart-line",
    "chart-pie",
    "table",
}


def emu_to_cm(emu: float) -> float:
    """将 PPTX 的 EMU 坐标单位转换为厘米。

    Args:
        emu: PPTX 内部使用的 EMU 数值。

    Returns:
        转换并保留两位小数后的厘米数值。
    """
    return round(emu / 360000.0, 2)


def get_shape_layout(shape) -> dict[str, float]:
    """读取 PPTX shape 在页面上的位置和尺寸。

    Args:
        shape: python-pptx 的 shape 对象。

    Returns:
        包含 `x`、`y`、`width`、`height` 的布局字典，单位为厘米。
    """
    return {
        "x": emu_to_cm(shape.left),
        "y": emu_to_cm(shape.top),
        "width": emu_to_cm(shape.width),
        "height": emu_to_cm(shape.height),
    }


def get_shape_kind(shape) -> str:
    """判断 PPTX shape 的粗粒度类型。

    Args:
        shape: python-pptx 的 shape 对象。

    Returns:
        shape 类型字符串。可能值包括 `text`、`table`、`chart-bar`、
        `chart-line`、`chart-pie`、`chart`、`other`。
    """
    if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
        return "table"

    if shape.shape_type == MSO_SHAPE_TYPE.CHART and hasattr(shape, "chart"):
        chart_type = shape.chart.chart_type
        if chart_type in _chart_type_values(
            "COLUMN_CLUSTERED",
            "COLUMN_STACKED",
            "COLUMN_STACKED_100",
            "BAR_CLUSTERED",
            "BAR_STACKED",
            "BAR_STACKED_100",
        ):
            return "chart-bar"
        if chart_type in _chart_type_values(
            "LINE",
            "LINE_MARKERS",
            "LINE_STACKED",
            "LINE_STACKED_100",
            "LINE_MARKERS_STACKED",
            "LINE_MARKERS_STACKED_100",
        ):
            return "chart-line"
        if chart_type in _chart_type_values("PIE", "PIE_EXPLODED", "DOUGHNUT"):
            return "chart-pie"
        return "chart"

    if shape.has_text_frame and shape.text.strip():
        return "text"

    return "other"


def extract_pptx_elements(pptx_path: Path, slide_idx: int = 0) -> dict[str, Any]:
    """从 PPTX 中抽取可编辑元素，不做语义判断。

    这一步只读取 PPTX 自身的结构信息，包括文本框、表格、图表和布局。
    不读取 benchmark gold 信息，也不抽取 chart/table 的二维数据。

    Args:
        pptx_path: 待解析的 PPTX 文件路径。
        slide_idx: 要解析的 slide 下标，当前流程默认只处理第 0 页。

    Returns:
        `observed_slide` 初始结构，包含 `slide_size` 和 `elements`。
        此时 `elements` 中还没有 `role` 字段。

    Raises:
        IndexError: 当 `slide_idx` 超出 PPTX 页数时抛出。
    """
    presentation = Presentation(pptx_path)
    if slide_idx >= len(presentation.slides):
        raise IndexError(f"Slide index {slide_idx} out of range.")

    slide = presentation.slides[slide_idx]
    elements: list[dict[str, Any]] = []

    for shape_idx, shape in enumerate(slide.shapes, 1):
        shape_kind = get_shape_kind(shape)
        if shape_kind == "other":
            continue

        element_id = str(shape_idx)
        element: dict[str, Any] = {
            "id": element_id,
            "layout": get_shape_layout(shape),
        }

        if shape_kind == "text":
            element.update(
                {
                    "type": "textBox",
                    "text": shape.text.strip(),
                }
            )
        elif shape_kind == "table":
            element["type"] = "table"
        else:
            element["type"] = "chart"
            element["_shape_kind"] = shape_kind

        elements.append(element)

    return {
        "slide_size": {
            "width": emu_to_cm(presentation.slide_width),
            "height": emu_to_cm(presentation.slide_height),
        },
        "elements": elements,
    }


def build_role_labeling_prompt(observed_slide: dict[str, Any]) -> str:
    """构造 VLM role 标注输入。

    这一步只组装 PPTX 已抽取出的元素列表、元素布局、文本内容和允许的
    role 集合。具体 role 标注规则存放在 `method/prompts/role_labeling_prompt.txt`。

    Args:
        observed_slide: `extract_pptx_elements` 输出的无 role slide 结构。

    Returns:
        传给 VLM user message 的输入字符串。
    """
    elements = []
    for element in observed_slide.get("elements", []):
        prompt_element = {
            "id": str(element.get("id")),
            "type": element.get("type"),
            "layout": element.get("layout"),
        }
        if element.get("type") == "textBox":
            prompt_element["text"] = element.get("text", "")
        else:
            prompt_element["shape_kind"] = element.get("_shape_kind") or element.get("type")
        elements.append(prompt_element)

    payload = {
        "slide_size": observed_slide.get("slide_size"),
        "elements": elements,
        "allowed_roles": sorted(ROLE_SET),
    }
    return f"Input elements:\n{json.dumps(payload, ensure_ascii=False, indent=2)}"


def _label_roles(
    client: Any,
    *,
    system_prompt: str,
    image_path: Path,
    observed_slide: dict[str, Any],
) -> list[dict[str, str]]:
    """调用 VLM 为 PPTX 元素标注语义 role。

    Args:
        client: 具备 `chat(...)` 方法的 OpenAI-compatible client。
        system_prompt: 从 prompt 文件读取的 role 标注规则。
        image_path: slide 渲染图路径，供 VLM 结合视觉布局判断 role。
        observed_slide: 无 role 的 PPTX 元素结构。

    Returns:
        VLM 返回的 role 标注列表，每项包含 `id` 和 `role`。

    Raises:
        ValueError: 当 VLM 返回 JSON 中缺少 `roles` 列表时抛出。
    """
    prompt = build_role_labeling_prompt(observed_slide)
    content = client.chat(
        system_prompt,
        prompt,
        image_path=image_path,
        response_format="json_object",
    )
    parsed = parse_json_object(content)
    roles = parsed.get("roles")
    if not isinstance(roles, list):
        raise ValueError(f"VLM role response must contain roles list: {content}")
    return roles


def validate_role_labels(
    observed_slide: dict[str, Any],
    roles: list[Any],
) -> list[dict[str, str]]:
    """校验 VLM role 标注是否完整且合法。

    Args:
        observed_slide: 无 role 的 PPTX 元素结构。
        roles: VLM 返回的 role 标注列表。

    Returns:
        规范化后的 role 标注列表，每项为 `{"id": ..., "role": ...}`。

    Raises:
        ValueError: 当 role 项格式错误、id 未知、id 重复、role 不在允许集合、
            或存在未标注元素时抛出。
    """
    expected_ids = {str(element.get("id")) for element in observed_slide.get("elements", [])}
    assignments: list[dict[str, str]] = []
    seen: set[str] = set()

    for item in roles:
        if not isinstance(item, dict):
            raise ValueError(f"Role assignment must be an object: {item}")
        element_id = str(item.get("id", "")).strip()
        role = str(item.get("role", "")).strip()
        if element_id not in expected_ids:
            raise ValueError(f"Unknown element id in role assignment: {element_id}")
        if element_id in seen:
            raise ValueError(f"Duplicate role assignment for element id: {element_id}")
        if role not in ROLE_SET:
            raise ValueError(f"Unsupported role '{role}' for element id {element_id}")
        assignments.append({"id": element_id, "role": role})
        seen.add(element_id)

    missing = expected_ids - seen
    if missing:
        raise ValueError(f"Missing role assignments for element ids: {sorted(missing)}")
    return assignments


def apply_role_labels(
    observed_slide: dict[str, Any],
    roles: list[dict[str, str]],
) -> dict[str, Any]:
    """把 role 标注合并回 PPTX 元素结构。

    Args:
        observed_slide: 无 role 的 PPTX 元素结构。
        roles: 已通过校验的 role 标注列表。

    Returns:
        带 `role` 字段的 `observed_slide`。内部临时字段如 `_shape_kind`
        会被移除，避免泄漏到后续表示中。
    """
    role_by_id = {item["id"]: item["role"] for item in roles}
    elements = []
    for element in observed_slide.get("elements", []):
        clean_element = {
            key: value for key, value in element.items() if not str(key).startswith("_")
        }
        clean_element["role"] = role_by_id[str(element.get("id"))]
        elements.append(clean_element)
    return {
        "slide_size": observed_slide.get("slide_size", {}),
        "elements": elements,
    }


class SlideParserAgent:
    """将 PPTX/PNG 解析为后续 verification 可读的 slide 表示。

    Parser 负责三件事：抽取 PPTX 元素、调用 VLM 标注 role、导出
    chart/table 数据 CSV，并组装 `ppt_representation`。
    """

    def __init__(
        self,
        *,
        model: str | None = None,
        api_key: str | None = None,
        base_url: str = "https://dashscope.aliyuncs.com/compatible-mode/v1",
        timeout_sec: int = 120,
        enable_thinking: bool | None = False,
        client: Any | None = None,
        role_labeling_prompt_path: Path = DEFAULT_ROLE_LABELING_PROMPT_PATH,
    ):
        """初始化 slide parser。

        Args:
            model: OpenAI-compatible VLM 模型名。未传入 `client` 时必须提供。
            api_key: API key；为空时由底层 `Client` 从环境变量读取。
            base_url: OpenAI-compatible API base URL。
            timeout_sec: 单次请求超时时间，单位秒。
            enable_thinking: 是否通过 `extra_body` 打开模型 thinking。
            client: 可选的已构造 client，主要用于测试或外部统一注入。
            role_labeling_prompt_path: role 标注 prompt 文件路径。

        Returns:
            None。

        Raises:
            ValueError: 当未提供 `client` 且 `model` 为空时抛出。
        """
        self.role_labeling_prompt = role_labeling_prompt_path.read_text(encoding="utf-8")
        if client is not None:
            self.client = client
        else:
            if model is None:
                raise ValueError("SlideParserAgent requires model when client is not provided.")
            self.client = Client(
                model=model,
                api_key=api_key,
                base_url=base_url,
                timeout_sec=timeout_sec,
                enable_thinking=enable_thinking,
            )

    def run(self, slide_input: SlideReviewInput) -> dict[str, Any]:
        """执行单页 slide 的 Phase 1 parser 流程。

        Args:
            slide_input: 当前 case 的输入路径，只包含 `pptx_path` 和
                `image_path`。

        Returns:
            包含两个顶层字段的字典：
            `observed_slide` 为带 role 的可编辑元素列表；
            `ppt_representation` 为简化后的 slide 表示，并引用导出的 CSV。
        """
        raw_elements = extract_pptx_elements(slide_input.pptx_path)
        role_labels = validate_role_labels(
            raw_elements,
            _label_roles(
                self.client,
                system_prompt=self.role_labeling_prompt,
                image_path=slide_input.image_path,
                observed_slide=raw_elements,
            ),
        )
        observed_slide = apply_role_labels(raw_elements, role_labels)
        ppt_representation = build_ppt_representation(
            pptx_path=slide_input.pptx_path,
            observed_slide=observed_slide,
        )
        return {
            "observed_slide": observed_slide,
            "ppt_representation": ppt_representation,
        }


def build_ppt_representation(
    *,
    pptx_path: Path,
    observed_slide: dict[str, Any],
    slide_idx: int = 0,
) -> dict[str, Any]:
    """把带 role 的元素整理为紧凑的 PPT 表示。

    这一步只做结构整理和 chart/table 数据导出，不做业务分析、不推断
    数据库字段、不判断 slide 是否有错误。

    Args:
        pptx_path: PPTX 文件路径，CSV 会写入该文件所在目录。
        observed_slide: 已合并 role 的 slide 元素结构。
        slide_idx: 要读取 chart/table 数据的 slide 下标。

    Returns:
        `ppt_representation` 字典，包含 `title`、`summary` 和
        `structured_tables`。每个 table body 通过 `data_path` 指向 CSV。
    """
    presentation = Presentation(pptx_path)
    slide = presentation.slides[slide_idx]
    elements = list(observed_slide.get("elements", []))
    title = _first_element_by_role(elements, "title")
    summary = _first_element_by_role(elements, "summary")
    captions = [element for element in elements if element.get("role") == "caption"]
    bodies = [
        element
        for element in elements
        if element.get("role") in {"chart-bar", "chart-line", "chart-pie", "table"}
    ]

    structured_tables = []
    output_dir = pptx_path.parent
    for table_idx, body_element in enumerate(bodies, 1):
        shape = slide.shapes[int(body_element["id"]) - 1]
        caption = _nearest_caption(body_element, captions)
        header, rows = _extract_body_table(shape=shape, role=str(body_element.get("role")))
        csv_path = output_dir / ("data.csv" if len(bodies) == 1 else f"data_{table_idx}.csv")
        _write_csv(csv_path, header, rows)
        structured_tables.append(
            {
                "caption": _compact_text_element(caption),
                "header": header,
                "body": {
                    "element_id": str(body_element.get("id")),
                    "type": body_element.get("role"),
                    "data_path": str(csv_path),
                },
            }
        )

    return {
        "title": _compact_text_element(title),
        "summary": _compact_text_element(summary),
        "structured_tables": structured_tables,
    }


def _first_element_by_role(elements: list[dict[str, Any]], role: str) -> dict[str, Any] | None:
    """按 role 查找第一个元素。

    Args:
        elements: 带 role 的元素列表。
        role: 要查找的 role 名称。

    Returns:
        第一个匹配元素；若不存在则返回 None。
    """
    for element in elements:
        if element.get("role") == role:
            return element
    return None


def _compact_text_element(element: dict[str, Any] | None) -> dict[str, Any] | None:
    """压缩文本元素，只保留后续需要的字段。

    Args:
        element: 带文本的元素；可以为 None。

    Returns:
        包含 `element_id` 和 `text` 的字典；输入为空时返回 None。
    """
    if not element:
        return None
    return {
        "element_id": str(element.get("id")),
        "text": str(element.get("text", "")).strip(),
    }


def _nearest_caption(
    body_element: dict[str, Any],
    captions: list[dict[str, Any]],
) -> dict[str, Any] | None:
    """为 chart/table body 选择空间位置最近的 caption。

    Args:
        body_element: role 为 chart/table 的 body 元素。
        captions: 候选 caption 元素列表。

    Returns:
        距离 body 中心点最近的 caption；没有候选时返回 None。
    """
    if not captions:
        return None
    body_center = _center(body_element.get("layout", {}))
    return min(
        captions,
        key=lambda caption: abs(body_center[0] - _center(caption.get("layout", {}))[0])
        + abs(body_center[1] - _center(caption.get("layout", {}))[1]),
    )


def _center(layout: dict[str, Any]) -> tuple[float, float]:
    """计算布局矩形中心点。

    Args:
        layout: 包含 `x`、`y`、`width`、`height` 的布局字典。

    Returns:
        `(center_x, center_y)` 中心点坐标。
    """
    return (
        float(layout.get("x", 0.0)) + float(layout.get("width", 0.0)) / 2.0,
        float(layout.get("y", 0.0)) + float(layout.get("height", 0.0)) / 2.0,
    )


def _extract_body_table(shape, role: str) -> tuple[list[str], list[list[Any]]]:
    """从 table 或 chart shape 中抽取二维数据。

    Args:
        shape: python-pptx 的 table/chart shape 对象。
        role: body 元素 role，用于区分 table 与 chart。

    Returns:
        二元组 `(header, rows)`。`header` 是 CSV 表头，`rows` 是数据行。
    """
    if role == "table":
        raw_rows = [
            [str(shape.table.cell(row_idx, col_idx).text).strip() for col_idx in range(len(shape.table.columns))]
            for row_idx in range(len(shape.table.rows))
        ]
        return (raw_rows[0] if raw_rows else []), raw_rows[1:]

    chart = shape.chart
    categories = [str(category.label) for category in chart.plots[0].categories]
    series = list(chart.series)
    header = ["category"] + [str(item.name) for item in series]
    rows = []
    for idx, category in enumerate(categories):
        rows.append([category] + [item.values[idx] if idx < len(item.values) else "" for item in series])
    return header, rows


def _write_csv(path: Path, header: list[str], rows: list[list[Any]]) -> None:
    """把二维数据写成 CSV。

    Args:
        path: 输出 CSV 路径。
        header: CSV 表头。
        rows: CSV 数据行。

    Returns:
        None。
    """
    with path.open("w", encoding="utf-8", newline="") as handle:
        writer = csv.writer(handle)
        writer.writerow(header)
        writer.writerows(rows)


def _chart_type_values(*names: str) -> set[Any]:
    """把 chart type 常量名转换为 python-pptx 枚举值集合。

    Args:
        *names: `XL_CHART_TYPE` 上的常量名。

    Returns:
        存在于当前 python-pptx 版本中的枚举值集合。
    """
    values = set()
    for name in names:
        if hasattr(XL_CHART_TYPE, name):
            values.add(getattr(XL_CHART_TYPE, name))
    return values
