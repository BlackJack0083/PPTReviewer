"""Content validation 的非 agent-tool 辅助函数。"""

from __future__ import annotations

import shutil
from pathlib import Path
from typing import Any

import pandas as pd
import yaml
from pptx import Presentation
from pptx.chart.data import ChartData


def write_content_artifacts(
    *,
    source_pptx: Path,
    analysis_state: dict[str, Any],
    artifact_dir: Path,
) -> dict[str, Any]:
    """写出 repaired PPTX、semantic YAML 和 table CSV。

    Args:
        source_pptx: 待修复的源 PPTX。
        analysis_state: content validation 后的最终 analysis state。
        artifact_dir: repaired YAML/CSV 文件写入目录。

    Returns:
        repaired PPTX、YAML 和 CSV 路径。
    """
    artifact_dir.mkdir(parents=True, exist_ok=True)
    data_paths = _write_table_data_files(
        analysis_state=analysis_state,
        artifact_dir=artifact_dir,
    )
    yaml_payload = _build_repaired_yaml(
        analysis_state=analysis_state,
        data_paths=data_paths,
    )
    yaml_path = artifact_dir / "repaired_slide.yaml"
    yaml_path.write_text(
        yaml.safe_dump(yaml_payload, sort_keys=False, allow_unicode=True),
        encoding="utf-8",
    )
    pptx_path = artifact_dir / "repaired_slide.pptx"
    _write_repaired_pptx(
        source_pptx=source_pptx,
        output_pptx=pptx_path,
        analysis_state=analysis_state,
    )
    return {
        "yaml_path": str(yaml_path),
        "data_paths": {str(key): value for key, value in data_paths.items()},
        "pptx_path": str(pptx_path),
    }


def _write_repaired_pptx(
    *,
    source_pptx: Path,
    output_pptx: Path,
    analysis_state: dict[str, Any],
) -> None:
    """把最终 analysis state 写回源 PPTX 的第一页。"""
    presentation = Presentation(source_pptx)
    slide = presentation.slides[0]

    text_elements = [
        analysis_state["title"],
        analysis_state["summary"],
        *(table["caption"] for table in analysis_state["tables"]),
    ]
    for element in text_elements:
        slide.shapes[int(element["element_id"]) - 1].text = element["text"]

    for table_state in analysis_state["tables"]:
        body = table_state["body"]
        shape = slide.shapes[int(body["element_id"]) - 1]
        dataframe = pd.read_csv(table_state["data_path"])
        if body["type"] == "table":
            _replace_table_data(shape.table, dataframe)
        else:
            _replace_chart_data(shape.chart, dataframe)

    presentation.save(output_pptx)
    Presentation(output_pptx)


def _replace_chart_data(chart: Any, dataframe: pd.DataFrame) -> None:
    """用 dataframe 第一列作为类别，其余列作为 series 替换 chart 数据。"""
    if dataframe.shape[1] < 2:
        raise ValueError("Chart data must contain one category and at least one series column.")
    chart_data = ChartData()
    chart_data.categories = dataframe.iloc[:, 0].tolist()
    for column in dataframe.columns[1:]:
        values = [None if pd.isna(value) else value for value in dataframe[column]]
        chart_data.add_series(str(column), values)
    chart.replace_data(chart_data)


def _replace_table_data(table: Any, dataframe: pd.DataFrame) -> None:
    """用 dataframe header 和 records 替换同尺寸 PPT table。"""
    rows = [list(dataframe.columns), *dataframe.fillna("").values.tolist()]
    expected_shape = (len(table.rows), len(table.columns))
    actual_shape = (len(rows), len(dataframe.columns))
    if actual_shape != expected_shape:
        raise ValueError(
            f"Repaired table shape {actual_shape} does not match PPT table "
            f"shape {expected_shape}."
        )
    for row_index, row in enumerate(rows):
        for column_index, value in enumerate(row):
            table.cell(row_index, column_index).text = str(value)


def _write_table_data_files(
    *,
    analysis_state: dict[str, Any],
    artifact_dir: Path,
) -> dict[int, str]:
    data_paths: dict[int, str] = {}
    for table_index, table_state in enumerate(analysis_state["tables"]):
        source_path = Path(table_state["data_path"])
        target_path = artifact_dir / f"table_{table_index}_repaired.csv"
        if source_path.resolve() != target_path.resolve():
            shutil.copyfile(source_path, target_path)
        data_paths[table_index] = str(target_path)
    return data_paths


def _build_repaired_yaml(
    *,
    analysis_state: dict[str, Any],
    data_paths: dict[int, str],
) -> dict[str, Any]:
    tables = []
    for table_index, table_state in enumerate(analysis_state.get("tables", [])):
        body = table_state.get("body") or {}
        tables.append(
            {
                "caption": table_state["caption"]["text"],
                "body": {
                    "presentation_type": body.get("type", ""),
                    "data_path": data_paths.get(
                        table_index,
                        table_state.get("data_path", ""),
                    ),
                },
                "data_source": table_state["caption"]["data_source"],
                "calculation_logic": table_state.get("calculation_logic", {}),
            }
        )
    return {
        "title": analysis_state["title"]["text"],
        "summary": analysis_state["summary"]["text"],
        "tables": tables,
    }
