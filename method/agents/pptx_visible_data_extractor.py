from __future__ import annotations

import re
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.enum.chart import XL_CHART_TYPE

YEAR_RE = re.compile(r"\b(20\d{2})\b")
MONTH_RE = re.compile(r"\b(20\d{2})-(0[1-9]|1[0-2])\b")
NUMBER_RE = re.compile(r"[-+]?\d[\d,]*(?:\.\d+)?")


def parse_number_token(text: str) -> float | None:
    match = NUMBER_RE.search(str(text))
    if not match:
        return None
    try:
        return float(match.group(0).replace(",", ""))
    except ValueError:
        return None


def infer_metric_kind(name: str) -> str:
    lower = name.lower()
    if "price" in lower:
        return "price"
    if "supply" in lower:
        return "supply"
    if "trade" in lower or "sales" in lower or "transaction" in lower:
        return "transaction"
    if "area" in lower:
        return "area"
    return "unknown"


def infer_category_granularity(categories: list[str]) -> str:
    if not categories:
        return "unknown"
    if all(MONTH_RE.search(value) for value in categories):
        return "month"
    if all(YEAR_RE.search(value) for value in categories):
        return "year"
    if any("m²" in value or "m2" in value.lower() or "sqm" in value.lower() for value in categories):
        return "area_segment"
    return "other"


def infer_chart_presentation_type(chart) -> str:
    chart_type = chart.chart_type
    if chart_type in {
        XL_CHART_TYPE.LINE,
        XL_CHART_TYPE.LINE_MARKERS,
        XL_CHART_TYPE.LINE_STACKED,
        XL_CHART_TYPE.LINE_STACKED_100,
        XL_CHART_TYPE.LINE_MARKERS_STACKED,
        XL_CHART_TYPE.LINE_MARKERS_STACKED_100,
    }:
        return "line chart"
    if chart_type in {XL_CHART_TYPE.PIE, XL_CHART_TYPE.PIE_EXPLODED, XL_CHART_TYPE.DOUGHNUT}:
        return "pie chart"
    return "bar chart"


def series_support_values(values: list[float]) -> list[float]:
    clean = [float(value) for value in values]
    if not clean:
        return []
    support = list(clean)
    support.extend(
        [
            clean[0],
            clean[-1],
            max(clean),
            min(clean),
            sum(clean) / len(clean),
            clean[-1] - clean[0],
        ]
    )
    return support


def infer_trend_direction(values: list[float]) -> str | None:
    if len(values) < 2:
        return None
    delta = values[-1] - values[0]
    if abs(delta) < 1e-9:
        return "flat"
    return "increase" if delta > 0 else "decrease"


class PPTXVisibleDataExtractor:
    """Extract visible chart/table payloads from slide.pptx without touching gold YAML."""

    def extract(self, pptx_path: Path, observed_slide: dict[str, Any]) -> dict[str, dict[str, Any]]:
        presentation = Presentation(pptx_path)
        slide = presentation.slides[0]
        extracted: dict[str, dict[str, Any]] = {}

        for element in observed_slide.get("elements", []):
            role = str(element.get("role", ""))
            if role not in {"chart-bar", "chart-line", "chart-pie", "table"}:
                continue
            element_id = str(element.get("id"))
            shape = slide.shapes[int(element_id) - 1]
            if role == "table":
                extracted[element_id] = self._extract_table(shape)
            else:
                extracted[element_id] = self._extract_chart(shape)
        return extracted

    def _extract_chart(self, shape) -> dict[str, Any]:
        chart = shape.chart
        categories = [str(category.label) for category in chart.plots[0].categories]
        series_payload = []
        metric_kinds = set()
        all_support_values: list[float] = []
        for series in chart.series:
            values = [float(value) for value in series.values]
            support_values = series_support_values(values)
            metric_kind = infer_metric_kind(series.name)
            metric_kinds.add(metric_kind)
            all_support_values.extend(support_values)
            series_payload.append(
                {
                    "name": str(series.name),
                    "metric_kind": metric_kind,
                    "values": values,
                    "support_values": support_values,
                    "trend_direction": infer_trend_direction(values),
                    "max_value": max(values) if values else None,
                }
            )

        years = sorted(
            {
                int(match.group(1))
                for category in categories
                for match in [YEAR_RE.search(category)]
                if match
            }
        )
        return {
            "body_kind": "chart",
            "presentation_type": infer_chart_presentation_type(chart),
            "categories": categories,
            "category_granularity": infer_category_granularity(categories),
            "time_range": [years[0], years[-1]] if years else None,
            "series": series_payload,
            "metric_kinds": sorted(kind for kind in metric_kinds if kind != "unknown"),
            "support_values": all_support_values,
        }

    def _extract_table(self, shape) -> dict[str, Any]:
        table = shape.table
        rows = [
            [str(table.cell(row_idx, col_idx).text).strip() for col_idx in range(len(table.columns))]
            for row_idx in range(len(table.rows))
        ]
        header = rows[0] if rows else []
        category_headers = header[1:] if len(header) > 1 else []
        row_payload = []
        metric_kinds = set()
        all_support_values: list[float] = []
        for row in rows[1:]:
            if not row:
                continue
            metric_name = row[0]
            values = [parse_number_token(value) for value in row[1:]]
            clean_values = [float(value) for value in values if value is not None]
            support_values = series_support_values(clean_values)
            metric_kind = infer_metric_kind(metric_name)
            metric_kinds.add(metric_kind)
            all_support_values.extend(support_values)
            row_payload.append(
                {
                    "name": metric_name,
                    "metric_kind": metric_kind,
                    "values": clean_values,
                    "support_values": support_values,
                    "trend_direction": infer_trend_direction(clean_values),
                    "max_value": max(clean_values) if clean_values else None,
                }
            )

        years = sorted(
            {
                int(match.group(1))
                for category in category_headers
                for match in [YEAR_RE.search(category)]
                if match
            }
        )
        return {
            "body_kind": "table",
            "presentation_type": "table",
            "categories": category_headers,
            "category_granularity": infer_category_granularity(category_headers),
            "time_range": [years[0], years[-1]] if years else None,
            "rows": row_payload,
            "metric_kinds": sorted(kind for kind in metric_kinds if kind != "unknown"),
            "support_values": all_support_values,
            "raw_rows": rows,
        }
