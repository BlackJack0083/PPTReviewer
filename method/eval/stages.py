from __future__ import annotations

import copy
from pathlib import Path
from typing import Any

import pandas as pd
from pptx import Presentation

from method.agents.content_validation.tools import execute_table_state
from method.agents.slide_parser.pptx_extract import extract_body_table, get_shape
from method.transformers import StatTransformer
from method.utils import read_dataframe_csv

NUMERIC_TOLERANCE = 0.5
ROLE_MAP = {"slide-title": "title", "body-text": "summary"}
CITY_TABLE_PREFIXES = {"Beijing": "beijing", "Guangzhou": "guangzhou", "Shenzhen": "shenzhen"}


def evaluate_parser(
    *,
    observed_slide: dict[str, Any],
    ppt_representation: dict[str, Any],
    injected_yaml: dict[str, Any],
    injected_yaml_path: Path,
) -> dict[str, Any]:
    """Require exact roles, caption-body pairing, and extracted body data."""
    elements = injected_yaml["template_slide"]["elements"]
    expected_roles = {
        str(element["id"]): ROLE_MAP.get(element["role"], element["role"])
        for element in elements
    }
    predicted_roles = {
        str(element["id"]): element["role"] for element in observed_slide["elements"]
    }
    roles_correct = predicted_roles == expected_roles

    expected_captions = [
        element for element in elements if element["role"] == "caption"
    ]
    expected_bodies = [
        element
        for element in elements
        if element["role"] in {"chart-bar", "chart-line", "chart-pie", "table"}
    ]
    predicted_tables = ppt_representation["structured_tables"]
    pairing_correct = len(predicted_tables) == len(expected_bodies) and all(
        str(predicted["caption"]["element_id"]) == str(caption["id"])
        and str(predicted["body"]["element_id"]) == str(body["id"])
        for predicted, caption, body in zip(
            predicted_tables,
            expected_captions,
            expected_bodies,
            strict=False,
        )
    )
    data_correct = len(predicted_tables) == len(expected_bodies) and all(
        _dataframes_equal(
            pd.read_csv(predicted["body"]["data_path"]),
            _yaml_body_dataframe(body, injected_yaml_path),
        )
        for predicted, body in zip(predicted_tables, expected_bodies, strict=False)
    )
    return {
        "success": roles_correct and pairing_correct and data_correct,
        "roles_correct": roles_correct,
        "pairing_correct": pairing_correct,
        "data_correct": data_correct,
    }


def evaluate_data_source_extraction(
    *,
    analysis_state: dict[str, Any],
    injected_yaml: dict[str, Any],
) -> dict[str, Any]:
    """Evaluate complete caption datasource objects, not individual slots."""
    expected = _caption_data_sources(injected_yaml)
    predicted = [
        _normalize_caption_data_source(table["caption"]["data_source"])
        for table in analysis_state["tables"]
    ]
    correct = sum(
        left == right for left, right in zip(predicted, expected, strict=False)
    )
    total = max(len(predicted), len(expected))
    return {
        "accuracy": correct / total if total else 0.0,
        "success": len(predicted) == len(expected) and correct == total,
        "correct": correct,
        "total": total,
    }


def evaluate_function_logic(
    *,
    analysis_state: dict[str, Any],
    ground_truth_yaml: dict[str, Any],
    ground_truth_yaml_path: Path,
) -> dict[str, Any]:
    """Execute predicted logic with GT datasource and compare displayed results."""
    expected_sources = _caption_data_sources(ground_truth_yaml)
    expected_bodies = [
        element
        for element in ground_truth_yaml["template_slide"]["elements"]
        if element["role"] in {"chart-bar", "chart-line", "chart-pie", "table"}
    ]
    correct = 0
    errors: list[str] = []
    for index, (table, source, body) in enumerate(
        zip(analysis_state["tables"], expected_sources, expected_bodies, strict=False)
    ):
        table_state = copy.deepcopy(table)
        table_state["caption"]["data_source"] = source
        try:
            computed = execute_table_state(
                table_state,
                transformer=StatTransformer(),
            )
            expected = _yaml_body_dataframe(body, ground_truth_yaml_path)
            if _dataframes_equal(computed, expected, ignore_first_column_name=True):
                correct += 1
            else:
                errors.append(f"table[{index}] result differs from GT")
        except (
            Exception
        ) as exc:  # Evaluation converts execution failure into an incorrect case.
            errors.append(f"table[{index}] {type(exc).__name__}: {exc}")
    total = max(len(analysis_state["tables"]), len(expected_bodies))
    return {
        "accuracy": correct / total if total else 0.0,
        "success": len(analysis_state["tables"]) == len(expected_bodies)
        and correct == total,
        "correct": correct,
        "total": total,
        "errors": errors,
    }


def evaluate_content_repair(
    *,
    corruption_record: dict[str, Any],
    repaired_pptx_path: Path,
    ground_truth_pptx_path: Path,
) -> dict[str, Any]:
    """Evaluate the repaired content for every injected operation."""
    operations = corruption_record["operations"]
    if not repaired_pptx_path.exists():
        return {
            "accuracy": 0.0,
            "success": False,
            "correct": 0,
            "total": len(operations),
        }
    repaired_slide = Presentation(repaired_pptx_path).slides[0]
    ground_truth_slide = Presentation(ground_truth_pptx_path).slides[0]
    correct = 0
    for operation in operations:
        shape_index = int(operation["element_id"]) - 1
        repaired_shape = repaired_slide.shapes[shape_index]
        ground_truth_shape = ground_truth_slide.shapes[shape_index]
        if operation["kind"] == "text_slot_edit":
            operation_correct = repaired_shape.text == ground_truth_shape.text
        else:
            operation_correct = _dataframes_equal(
                _shape_dataframe(repaired_shape, operation["role"]),
                _shape_dataframe(ground_truth_shape, operation["role"]),
            )
        correct += int(operation_correct)
    total = len(operations)
    return {
        "accuracy": correct / total if total else 0.0,
        "success": correct == total,
        "correct": correct,
        "total": total,
    }


def compare_presentations(repaired_path: Path, ground_truth_path: Path) -> bool:
    """Return whether repaired PPTX exactly matches GT semantics and layout."""
    if not repaired_path.exists():
        return False
    try:
        repaired = Presentation(repaired_path)
        expected = Presentation(ground_truth_path)
    except Exception:
        return False
    if len(repaired.slides) != len(expected.slides):
        return False
    for repaired_slide, expected_slide in zip(
        repaired.slides,
        expected.slides,
        strict=True,
    ):
        if len(repaired_slide.shapes) != len(expected_slide.shapes):
            return False
        for repaired_shape, expected_shape in zip(
            repaired_slide.shapes,
            expected_slide.shapes,
            strict=True,
        ):
            if (
                repaired_shape.shape_type != expected_shape.shape_type
                or repaired_shape.left != expected_shape.left
                or repaired_shape.top != expected_shape.top
                or repaired_shape.width != expected_shape.width
                or repaired_shape.height != expected_shape.height
            ):
                return False
            if (
                repaired_shape.has_text_frame
                and repaired_shape.text != expected_shape.text
            ):
                return False
            shape_type = get_shape(expected_shape)
            if shape_type in {"chart-bar", "chart-line", "chart-pie", "table"} and not (
                _dataframes_equal(
                    _shape_dataframe(repaired_shape, shape_type),
                    _shape_dataframe(expected_shape, shape_type),
                )
            ):
                return False
    return True


def _caption_data_sources(yaml_data: dict[str, Any]) -> list[dict[str, Any]]:
    captions = [
        element
        for element in yaml_data["template_slide"]["elements"]
        if element.get("text_binding", {}).get("kind") == "caption"
    ]
    sources = []
    for caption, slide_filter in zip(
        captions,
        yaml_data["slide_filters"],
        strict=True,
    ):
        slots = caption["text_binding"]["slots"]
        city = str(slots["Geo_City_Name"]["value"])
        start_year = str(slots["Temporal_Start_Year"]["value"])
        end_year = str(slots["Temporal_End_Year"]["value"])
        table = _caption_table_name(city, slide_filter)
        sources.append(
            {
                "connection": {"table": table},
                "select_columns": sorted(slide_filter["select_columns"]),
                "filters": {
                    "city": city,
                    "block": str(slots["Geo_Block_Name"]["value"]),
                    "start_date": f"{start_year}-01-01" if start_year else "",
                    "end_date": f"{end_year}-12-31" if end_year else "",
                },
            }
        )
    return sources


def _normalize_caption_data_source(value: dict[str, Any]) -> dict[str, Any]:
    normalized = _normalize_final_data_source(value)
    normalized["select_columns"] = sorted(value["select_columns"])
    return normalized


def _normalize_final_data_source(value: dict[str, Any]) -> dict[str, Any]:
    return {
        "connection": {"table": str(value["connection"]["table"]).lower()},
        "filters": {
            field: str(value["filters"][field])
            for field in ("city", "block", "start_date", "end_date")
        },
    }


def _table_name(value: str | list[str]) -> str:
    return str(value[0] if isinstance(value, list) else value).lower()


def _caption_table_name(city: str, slide_filter: dict[str, Any]) -> str:
    if not city:
        return ""
    base_table = _table_name(slide_filter["connection"]["table"])
    suffix = "resale_house" if base_table.endswith("resale_house") else "new_house"
    return f"{CITY_TABLE_PREFIXES.get(city, base_table.rsplit('_', 2)[0])}_{suffix}"


def _yaml_body_dataframe(element: dict[str, Any], yaml_path: Path) -> pd.DataFrame:
    dataframe = read_dataframe_csv(yaml_path.parent / element["data"])
    if element["role"] == "table":
        return dataframe.reset_index(drop=True)
    transposed = dataframe.T.reset_index()
    transposed.columns = ["category", *[str(index) for index in dataframe.index]]
    return transposed


def _shape_dataframe(shape: Any, role: str) -> pd.DataFrame:
    header, rows = extract_body_table(shape, role)
    return pd.DataFrame(rows, columns=header)


def _dataframes_equal(
    left: pd.DataFrame,
    right: pd.DataFrame,
    *,
    ignore_first_column_name: bool = False,
) -> bool:
    left = left.reset_index(drop=True)
    right = right.reset_index(drop=True)
    if left.shape != right.shape:
        return False
    left_columns = [str(column).strip() for column in left.columns]
    right_columns = [str(column).strip() for column in right.columns]
    if ignore_first_column_name:
        left_columns[0] = right_columns[0]
    if left_columns != right_columns:
        return False
    for left_value, right_value in zip(
        left.to_numpy().flat,
        right.to_numpy().flat,
        strict=True,
    ):
        if pd.isna(left_value) and pd.isna(right_value):
            continue
        try:
            if abs(float(left_value) - float(right_value)) <= NUMERIC_TOLERANCE:
                continue
        except (TypeError, ValueError):
            pass
        if str(left_value).strip() != str(right_value).strip():
            return False
    return True
