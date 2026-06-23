from __future__ import annotations

import copy
from collections import Counter
from collections.abc import Callable
from pathlib import Path
from typing import Any

import pandas as pd
import yaml
from pptx import Presentation

from engine.data_files import read_dataframe_csv
from method.agents.content_validation.tools import execute_table_state
from method.agents.slide_parser.agent import _extract_body_table, get_shape
from method.transformers import StatTransformer

NUMERIC_TOLERANCE = 0.5
ROLE_MAP = {"slide-title": "title", "body-text": "summary"}


class SlideReviewEvaluator:
    """Evaluate detection, stage outputs, and repaired PPTX against benchmark GT."""

    def evaluate_case(
        self,
        *,
        result: dict[str, Any],
        corruption_record: dict[str, Any],
        injected_yaml_path: Path,
        ground_truth_yaml_path: Path,
        ground_truth_pptx_path: Path,
    ) -> dict[str, Any]:
        """Evaluate one completed slide-review case."""
        injected_yaml = _load_yaml(injected_yaml_path)
        ground_truth_yaml = _load_yaml(ground_truth_yaml_path)
        repaired_pptx_path = Path(result["repaired_artifacts"]["pptx_path"])

        detection = evaluate_detection(result["detected_issues"], corruption_record)
        parser = evaluate_parser(
            observed_slide=result["observed_slide"],
            ppt_representation=result["ppt_representation"],
            injected_yaml=injected_yaml,
            injected_yaml_path=injected_yaml_path,
        )
        data_source_extraction = evaluate_data_source_extraction(
            analysis_state=result["slide_analysis_state"],
            injected_yaml=injected_yaml,
        )
        function_logic = evaluate_function_logic(
            analysis_state=result["slide_analysis_state"],
            ground_truth_yaml=ground_truth_yaml,
            ground_truth_yaml_path=ground_truth_yaml_path,
        )
        data_source_validation = {
            "success": _normalize_final_data_source(
                result["analysis_state"]["final_data_source"]
            )
            == _ground_truth_data_source(ground_truth_yaml)
        }
        content_repair = evaluate_content_repair(
            corruption_record=corruption_record,
            repaired_pptx_path=repaired_pptx_path,
            ground_truth_pptx_path=ground_truth_pptx_path,
        )
        task_success = compare_presentations(
            repaired_pptx_path,
            ground_truth_pptx_path,
        )
        return {
            "detection": detection,
            "task_success": task_success,
            "stages": {
                "parser": parser,
                "data_source_extraction": data_source_extraction,
                "function_logic": function_logic,
                "data_source_validation": data_source_validation,
                "content_repair": content_repair,
            },
        }

    def evaluate_partial_case(
        self,
        *,
        partial_result: dict[str, Any],
        corruption_record: dict[str, Any],
        injected_yaml_path: Path,
        ground_truth_yaml_path: Path,
    ) -> dict[str, Any]:
        """Evaluate stages that completed before a workflow failure."""
        injected_yaml = _load_yaml(injected_yaml_path)
        ground_truth_yaml = _load_yaml(ground_truth_yaml_path)

        stages: dict[str, Any] = {
            "parser": {"success": False},
            "data_source_extraction": {"success": False},
            "function_logic": {"success": False},
            "data_source_validation": {"success": False},
            "content_repair": _content_repair_failure(corruption_record),
        }
        if (
            "observed_slide" in partial_result
            and "ppt_representation" in partial_result
        ):
            stages["parser"] = evaluate_parser(
                observed_slide=partial_result["observed_slide"],
                ppt_representation=partial_result["ppt_representation"],
                injected_yaml=injected_yaml,
                injected_yaml_path=injected_yaml_path,
            )
        if "slide_analysis_state" in partial_result:
            stages["data_source_extraction"] = evaluate_data_source_extraction(
                analysis_state=partial_result["slide_analysis_state"],
                injected_yaml=injected_yaml,
            )
            stages["function_logic"] = evaluate_function_logic(
                analysis_state=partial_result["slide_analysis_state"],
                ground_truth_yaml=ground_truth_yaml,
                ground_truth_yaml_path=ground_truth_yaml_path,
            )
        if "analysis_state" in partial_result:
            stages["data_source_validation"] = {
                "success": _normalize_final_data_source(
                    partial_result["analysis_state"]["final_data_source"]
                )
                == _ground_truth_data_source(ground_truth_yaml)
            }

        return {
            "detection": evaluate_detection(
                partial_result.get("detected_issues", []),
                corruption_record,
            ),
            "task_success": False,
            "stages": stages,
        }


def failure_metrics(corruption_record: dict[str, Any]) -> dict[str, Any]:
    """Return zero-valued metrics for a pipeline failure."""
    return {
        "detection": evaluate_detection([], corruption_record),
        "task_success": False,
        "stages": {
            "parser": {"success": False},
            "data_source_extraction": {"success": False},
            "function_logic": {"success": False},
            "data_source_validation": {"success": False},
            "content_repair": _content_repair_failure(corruption_record),
        },
    }


def _content_repair_failure(corruption_record: dict[str, Any]) -> dict[str, Any]:
    operation_count = len(corruption_record["operations"])
    return {
        "accuracy": 0.0,
        "success": False,
        "correct": 0,
        "total": operation_count,
    }


def aggregate_metrics(metrics: list[dict[str, Any]]) -> dict[str, Any]:
    """Aggregate the three primary metrics and five stage accuracies."""
    error_type_f1 = _aggregate_f1_by_label(
        case["detection"]["error_type"] for case in metrics
    )
    issue_f1 = _aggregate_f1_by_label(case["detection"]["issue"] for case in metrics)

    case_count = len(metrics)
    content_correct = sum(
        case["stages"]["content_repair"]["correct"] for case in metrics
    )
    content_total = sum(case["stages"]["content_repair"]["total"] for case in metrics)
    return {
        "error_type_macro_f1": error_type_f1["macro_f1"],
        "error_type_exact_accuracy": _exact_accuracy(metrics, "error_type"),
        "issue_macro_f1": issue_f1["macro_f1"],
        "issue_exact_accuracy": _exact_accuracy(metrics, "issue"),
        "task_success_rate": (
            sum(case["task_success"] for case in metrics) / case_count
            if case_count
            else 0.0
        ),
        "stage_accuracy": {
            stage: (
                sum(case["stages"][stage]["success"] for case in metrics) / case_count
                if case_count
                else 0.0
            )
            for stage in (
                "parser",
                "data_source_extraction",
                "function_logic",
                "data_source_validation",
            )
        }
        | {"content_repair": content_correct / content_total if content_total else 0.0},
        "error_type_f1_by_label": error_type_f1["f1_by_label"],
        "issue_f1_by_label": issue_f1["f1_by_label"],
    }


def evaluate_detection(
    detected_issues: list[dict[str, Any]],
    corruption_record: dict[str, Any],
) -> dict[str, Any]:
    """Evaluate coarse error type detection and fine-grained issue detection."""
    gold_items = [
        {**operation, "error_type": error_type}
        for operation in corruption_record["operations"]
        for error_type in operation["error_types"]
    ]
    return {
        "error_type": _label_metrics(
            predicted=Counter(_error_type_label(issue) for issue in detected_issues),
            gold=Counter(_error_type_label(item) for item in gold_items),
            label_to_dict=_error_type_label_to_dict,
        ),
        "issue": _label_metrics(
            predicted=Counter(_issue_label(issue) for issue in detected_issues),
            gold=Counter(_issue_label(item) for item in gold_items),
            label_to_dict=_issue_label_to_dict,
        ),
    }


def _label_metrics(
    *,
    predicted: Counter[tuple[str, ...]],
    gold: Counter[tuple[str, ...]],
    label_to_dict: Callable[[tuple[str, ...]], dict[str, str]],
) -> dict[str, Any]:
    labels = sorted(set(predicted) | set(gold))
    tp = Counter({label: min(predicted[label], gold[label]) for label in labels})
    fp = Counter({label: max(predicted[label] - gold[label], 0) for label in labels})
    fn = Counter({label: max(gold[label] - predicted[label], 0) for label in labels})
    tp_total = sum(tp.values())
    fp_total = sum(fp.values())
    fn_total = sum(fn.values())
    precision = tp_total / (tp_total + fp_total) if tp_total + fp_total else 0.0
    recall = tp_total / (tp_total + fn_total) if tp_total + fn_total else 0.0
    f1 = 2 * precision * recall / (precision + recall) if precision + recall else 0.0
    f1_by_label = [
        {
            "label": label_to_dict(label),
            "precision": (
                tp[label] / (tp[label] + fp[label]) if tp[label] + fp[label] else 0.0
            ),
            "recall": (
                tp[label] / (tp[label] + fn[label]) if tp[label] + fn[label] else 0.0
            ),
        }
        for label in labels
    ]
    for item in f1_by_label:
        item["f1"] = (
            2
            * item["precision"]
            * item["recall"]
            / (item["precision"] + item["recall"])
            if item["precision"] + item["recall"]
            else 0.0
        )
    return {
        "precision": precision,
        "recall": recall,
        "f1": f1,
        "macro_f1": (
            sum(item["f1"] for item in f1_by_label) / len(f1_by_label)
            if f1_by_label
            else 0.0
        ),
        "exact_match": predicted == gold,
        "predicted": _counter_payload(predicted, label_to_dict),
        "gold": _counter_payload(gold, label_to_dict),
        "tp": _counter_payload(tp, label_to_dict),
        "fp": _counter_payload(fp, label_to_dict),
        "fn": _counter_payload(fn, label_to_dict),
        "f1_by_label": f1_by_label,
    }


def _counter_payload(
    counter: Counter[tuple[str, ...]],
    label_to_dict: Callable[[tuple[str, ...]], dict[str, str]],
) -> list[dict[str, Any]]:
    payload = []
    for label, count in sorted(counter.items()):
        if count <= 0:
            continue
        item: dict[str, Any] = label_to_dict(label)
        if count > 1:
            item["count"] = count
        payload.append(item)
    return payload


def _aggregate_f1_by_label(
    metric_blocks: Any,
) -> dict[str, Any]:
    predicted = Counter()
    gold = Counter()
    for block in metric_blocks:
        predicted.update(_payload_counter(block["predicted"]))
        gold.update(_payload_counter(block["gold"]))
    metrics = _label_metrics(
        predicted=predicted,
        gold=gold,
        label_to_dict=_label_payload,
    )
    return {
        "macro_f1": metrics["macro_f1"],
        "f1_by_label": metrics["f1_by_label"],
    }


def _payload_counter(items: list[dict[str, Any]]) -> Counter[tuple[str, ...]]:
    counter: Counter[tuple[str, ...]] = Counter()
    for item in items:
        count = int(item.get("count", 1))
        if "scope_error_type" in item:
            label = (
                str(item["error_type"]),
                str(item["scope_error_type"]),
                str(item["field"]),
            )
        elif set(item) <= {"error_type", "count"}:
            label = (str(item["error_type"]),)
        else:
            label = (str(item["error_type"]), str(item["target"]))
        counter[label] += count
    return counter


def _exact_accuracy(metrics: list[dict[str, Any]], level: str) -> float:
    if not metrics:
        return 0.0
    return sum(case["detection"][level]["exact_match"] for case in metrics) / len(
        metrics
    )


def _error_type_label(item: dict[str, Any]) -> tuple[str, ...]:
    return (item["error_type"],)


def _issue_label(item: dict[str, Any]) -> tuple[str, ...]:
    error_type = item["error_type"]
    if error_type == "scope_error":
        return (error_type, item["scope_error_type"], item["field"])
    return (error_type, item["target"])


def _error_type_label_to_dict(label: tuple[str, ...]) -> dict[str, str]:
    return {"error_type": label[0]}


def _issue_label_to_dict(label: tuple[str, ...]) -> dict[str, str]:
    if label[0] == "scope_error":
        return {
            "error_type": label[0],
            "scope_error_type": label[1],
            "field": label[2],
        }
    return {"error_type": label[0], "target": label[1]}


def _label_payload(label: tuple[str, ...]) -> dict[str, str]:
    if len(label) == 1:
        return _error_type_label_to_dict(label)
    return _issue_label_to_dict(label)


def evaluate_parser(
    *,
    observed_slide: dict[str, Any],
    ppt_representation: dict[str, Any],
    injected_yaml: dict[str, Any],
    injected_yaml_path: Path,
) -> dict[str, Any]:
    """Require exact roles, caption-body pairing, and extracted body data."""
    elements = injected_yaml["template_slide"]["elements"]
    expected_roles = {
        str(element["id"]): ROLE_MAP.get(element["role"], element["role"])
        for element in elements
    }
    predicted_roles = {
        str(element["id"]): element["role"] for element in observed_slide["elements"]
    }
    roles_correct = predicted_roles == expected_roles

    expected_captions = [
        element for element in elements if element["role"] == "caption"
    ]
    expected_bodies = [
        element
        for element in elements
        if element["role"] in {"chart-bar", "chart-line", "chart-pie", "table"}
    ]
    predicted_tables = ppt_representation["structured_tables"]
    pairing_correct = len(predicted_tables) == len(expected_bodies) and all(
        str(predicted["caption"]["element_id"]) == str(caption["id"])
        and str(predicted["body"]["element_id"]) == str(body["id"])
        for predicted, caption, body in zip(
            predicted_tables,
            expected_captions,
            expected_bodies,
            strict=False,
        )
    )
    data_correct = len(predicted_tables) == len(expected_bodies) and all(
        _dataframes_equal(
            pd.read_csv(predicted["body"]["data_path"]),
            _yaml_body_dataframe(body, injected_yaml_path),
        )
        for predicted, body in zip(predicted_tables, expected_bodies, strict=False)
    )
    return {
        "success": roles_correct and pairing_correct and data_correct,
        "roles_correct": roles_correct,
        "pairing_correct": pairing_correct,
        "data_correct": data_correct,
    }


def evaluate_data_source_extraction(
    *,
    analysis_state: dict[str, Any],
    injected_yaml: dict[str, Any],
) -> dict[str, Any]:
    """Evaluate complete caption datasource objects, not individual slots."""
    expected = _caption_data_sources(injected_yaml)
    predicted = [
        _normalize_caption_data_source(table["caption"]["data_source"])
        for table in analysis_state["tables"]
    ]
    correct = sum(
        left == right for left, right in zip(predicted, expected, strict=False)
    )
    total = max(len(predicted), len(expected))
    return {
        "accuracy": correct / total if total else 0.0,
        "success": len(predicted) == len(expected) and correct == total,
        "correct": correct,
        "total": total,
    }


def evaluate_function_logic(
    *,
    analysis_state: dict[str, Any],
    ground_truth_yaml: dict[str, Any],
    ground_truth_yaml_path: Path,
) -> dict[str, Any]:
    """Execute predicted logic with GT datasource and compare displayed results."""
    expected_sources = _caption_data_sources(ground_truth_yaml)
    expected_bodies = [
        element
        for element in ground_truth_yaml["template_slide"]["elements"]
        if element["role"] in {"chart-bar", "chart-line", "chart-pie", "table"}
    ]
    correct = 0
    errors: list[str] = []
    for index, (table, source, body) in enumerate(
        zip(analysis_state["tables"], expected_sources, expected_bodies, strict=False)
    ):
        table_state = copy.deepcopy(table)
        table_state["caption"]["data_source"] = source
        try:
            computed = execute_table_state(
                table_state,
                transformer=StatTransformer(),
            )
            expected = _yaml_body_dataframe(body, ground_truth_yaml_path)
            if _dataframes_equal(computed, expected, ignore_first_column_name=True):
                correct += 1
            else:
                errors.append(f"table[{index}] result differs from GT")
        except (
            Exception
        ) as exc:  # Evaluation converts execution failure into an incorrect case.
            errors.append(f"table[{index}] {type(exc).__name__}: {exc}")
    total = max(len(analysis_state["tables"]), len(expected_bodies))
    return {
        "accuracy": correct / total if total else 0.0,
        "success": len(analysis_state["tables"]) == len(expected_bodies)
        and correct == total,
        "correct": correct,
        "total": total,
        "errors": errors,
    }


def evaluate_content_repair(
    *,
    corruption_record: dict[str, Any],
    repaired_pptx_path: Path,
    ground_truth_pptx_path: Path,
) -> dict[str, Any]:
    """Evaluate the repaired content for every injected operation."""
    operations = corruption_record["operations"]
    if not repaired_pptx_path.exists():
        return {
            "accuracy": 0.0,
            "success": False,
            "correct": 0,
            "total": len(operations),
        }
    repaired_slide = Presentation(repaired_pptx_path).slides[0]
    ground_truth_slide = Presentation(ground_truth_pptx_path).slides[0]
    correct = 0
    for operation in operations:
        shape_index = int(operation["element_id"]) - 1
        repaired_shape = repaired_slide.shapes[shape_index]
        ground_truth_shape = ground_truth_slide.shapes[shape_index]
        if operation["kind"] == "text_slot_edit":
            repaired = repaired_shape.text
            expected = ground_truth_shape.text
            operation_correct = repaired == expected
        else:
            role = operation["role"]
            operation_correct = _dataframes_equal(
                _shape_dataframe(repaired_shape, role),
                _shape_dataframe(ground_truth_shape, role),
            )
        correct += int(operation_correct)
    total = len(operations)
    return {
        "accuracy": correct / total if total else 0.0,
        "success": correct == total,
        "correct": correct,
        "total": total,
    }


def compare_presentations(repaired_path: Path, ground_truth_path: Path) -> bool:
    """Return whether repaired PPTX exactly matches GT semantics and layout."""
    if not repaired_path.exists():
        return False
    try:
        repaired = Presentation(repaired_path)
        expected = Presentation(ground_truth_path)
    except Exception:
        return False
    if len(repaired.slides) != len(expected.slides):
        return False
    for repaired_slide, expected_slide in zip(
        repaired.slides, expected.slides, strict=True
    ):
        if len(repaired_slide.shapes) != len(expected_slide.shapes):
            return False
        for repaired_shape, expected_shape in zip(
            repaired_slide.shapes,
            expected_slide.shapes,
            strict=True,
        ):
            if (
                repaired_shape.shape_type != expected_shape.shape_type
                or repaired_shape.left != expected_shape.left
                or repaired_shape.top != expected_shape.top
                or repaired_shape.width != expected_shape.width
                or repaired_shape.height != expected_shape.height
            ):
                return False
            if (
                repaired_shape.has_text_frame
                and repaired_shape.text != expected_shape.text
            ):
                return False
            shape_type = get_shape(expected_shape)
            if shape_type in {"chart-bar", "chart-line", "chart-pie", "table"} and not (
                _dataframes_equal(
                    _shape_dataframe(repaired_shape, shape_type),
                    _shape_dataframe(expected_shape, shape_type),
                )
            ):
                return False
    return True


def _caption_data_sources(yaml_data: dict[str, Any]) -> list[dict[str, Any]]:
    captions = [
        element
        for element in yaml_data["template_slide"]["elements"]
        if element.get("text_binding", {}).get("kind") == "caption"
    ]
    bodies = [
        element
        for element in yaml_data["template_slide"]["elements"]
        if element["role"] in {"chart-bar", "chart-line", "chart-pie", "table"}
    ]
    sources = []
    for caption, body, slide_filter in zip(
        captions,
        bodies,
        yaml_data["slide_filters"],
        strict=True,
    ):
        slots = caption["text_binding"]["slots"]
        city = str(slots["Geo_City_Name"]["value"])
        start_year = str(slots["Temporal_Start_Year"]["value"])
        end_year = str(slots["Temporal_End_Year"]["value"])
        table = _table_name(slide_filter["connection"]["table"]) if city else ""
        sources.append(
            {
                "connection": {"table": table},
                "select_columns": sorted(_required_columns(slide_filter, body)),
                "filters": {
                    "city": city,
                    "block": str(slots["Geo_Block_Name"]["value"]),
                    "start_date": f"{start_year}-01-01" if start_year else "",
                    "end_date": f"{end_year}-12-31" if end_year else "",
                },
            }
        )
    return sources


def _required_columns(
    slide_filter: dict[str, Any],
    body: dict[str, Any],
) -> set[str]:
    columns = set(slide_filter["select_columns"])
    logic = body["args"]
    for dimension in logic.get("dimensions", []):
        columns.add(dimension["source_col"])
    for metric in logic.get("metrics", []):
        columns.add(metric["source_col"])
        columns.update(metric.get("filter_condition", {}))
    return columns


def _normalize_caption_data_source(value: dict[str, Any]) -> dict[str, Any]:
    normalized = _normalize_final_data_source(value)
    normalized["select_columns"] = sorted(value["select_columns"])
    return normalized


def _normalize_final_data_source(value: dict[str, Any]) -> dict[str, Any]:
    return {
        "connection": {"table": str(value["connection"]["table"]).lower()},
        "filters": {
            field: str(value["filters"][field])
            for field in ("city", "block", "start_date", "end_date")
        },
    }


def _ground_truth_data_source(yaml_data: dict[str, Any]) -> dict[str, Any]:
    filters = yaml_data["query_filters"]
    return {
        "connection": {
            "table": _table_name(yaml_data["slide_filters"][0]["connection"]["table"])
        },
        "filters": {
            field: str(filters[field])
            for field in ("city", "block", "start_date", "end_date")
        },
    }


def _table_name(value: str | list[str]) -> str:
    return str(value[0] if isinstance(value, list) else value).lower()


def _yaml_body_dataframe(element: dict[str, Any], yaml_path: Path) -> pd.DataFrame:
    dataframe = read_dataframe_csv(yaml_path.parent / element["data"])
    if element["role"] == "table":
        return dataframe.reset_index(drop=True)
    transposed = dataframe.T.reset_index()
    transposed.columns = ["category", *[str(index) for index in dataframe.index]]
    return transposed


def _shape_dataframe(shape: Any, role: str) -> pd.DataFrame:
    header, rows = _extract_body_table(shape, role)
    return pd.DataFrame(rows, columns=header)


def _dataframes_equal(
    left: pd.DataFrame,
    right: pd.DataFrame,
    *,
    ignore_first_column_name: bool = False,
) -> bool:
    left = left.reset_index(drop=True)
    right = right.reset_index(drop=True)
    if left.shape != right.shape:
        return False
    left_columns = [str(column).strip() for column in left.columns]
    right_columns = [str(column).strip() for column in right.columns]
    if ignore_first_column_name:
        left_columns[0] = right_columns[0]
    if left_columns != right_columns:
        return False
    for left_value, right_value in zip(
        left.to_numpy().flat,
        right.to_numpy().flat,
        strict=True,
    ):
        if pd.isna(left_value) and pd.isna(right_value):
            continue
        try:
            if abs(float(left_value) - float(right_value)) <= NUMERIC_TOLERANCE:
                continue
        except (TypeError, ValueError):
            pass
        if str(left_value).strip() != str(right_value).strip():
            return False
    return True


def _load_yaml(path: Path) -> dict[str, Any]:
    value = yaml.safe_load(path.read_text(encoding="utf-8"))
    if not isinstance(value, dict):
        raise ValueError(f"Expected YAML object: {path}")
    return value
