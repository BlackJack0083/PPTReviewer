from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_SHAPE_TYPE


def emu_to_cm(emu: float) -> float:
    return round(emu / 360000.0, 2)


def get_shape_layout(shape) -> dict[str, float]:
    return {
        "x": emu_to_cm(shape.left),
        "y": emu_to_cm(shape.top),
        "width": emu_to_cm(shape.width),
        "height": emu_to_cm(shape.height),
    }


def extract_pptx_elements(pptx_path: Path, slide_idx: int = 0) -> dict:
    """Extract editable PPTX elements without semantic roles or chart/table data."""
    presentation = Presentation(pptx_path)
    if slide_idx >= len(presentation.slides):
        raise IndexError(f"Slide index {slide_idx} out of range.")

    slide = presentation.slides[slide_idx]
    elements: list[dict] = []

    for shape_idx, shape in enumerate(slide.shapes, 1):
        shape_kind = get_shape_kind(shape)
        if shape_kind == "other":
            continue

        element_id = str(shape_idx)
        element: dict = {
            "id": element_id,
            "layout": get_shape_layout(shape),
        }

        if shape_kind == "text":
            element.update(
                {
                    "type": "textBox",
                    "text": shape.text.strip(),
                }
            )
        elif shape_kind == "table":
            element["type"] = "table"
        else:
            element["type"] = "chart"
            element["_shape_kind"] = shape_kind

        elements.append(element)

    return {
        "slide_size": {
            "width": emu_to_cm(presentation.slide_width),
            "height": emu_to_cm(presentation.slide_height),
        },
        "elements": elements,
    }


def get_shape_kind(shape) -> str:
    if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
        return "table"

    if shape.shape_type == MSO_SHAPE_TYPE.CHART and hasattr(shape, "chart"):
        chart_type = shape.chart.chart_type
        if chart_type in _chart_type_values(
            "COLUMN_CLUSTERED",
            "COLUMN_STACKED",
            "COLUMN_STACKED_100",
            "BAR_CLUSTERED",
            "BAR_STACKED",
            "BAR_STACKED_100",
        ):
            return "chart-bar"
        if chart_type in _chart_type_values(
            "LINE",
            "LINE_MARKERS",
            "LINE_STACKED",
            "LINE_STACKED_100",
            "LINE_MARKERS_STACKED",
            "LINE_MARKERS_STACKED_100",
        ):
            return "chart-line"
        if chart_type in _chart_type_values("PIE", "PIE_EXPLODED", "DOUGHNUT"):
            return "chart-pie"
        return "chart"

    if shape.has_text_frame and shape.text.strip():
        return "text"

    return "other"


def _chart_type_values(*names: str) -> set:
    values = set()
    for name in names:
        if hasattr(XL_CHART_TYPE, name):
            values.add(getattr(XL_CHART_TYPE, name))
    return values
