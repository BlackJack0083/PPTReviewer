# ppt_operation.py
import os
from pathlib import Path
from typing import Self

import pandas as pd
import pptx_ea_font
from loguru import logger
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_DATA_LABEL_POSITION, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.slide import Slide
from pptx.util import Cm, Pt

# 引入 Pydantic 模型
from ppt_schemas import (
    ChartConfigModel,
    LayoutModel,
    RectangleStyleModel,
    TextContentModel,
)

CHART_COLORS_THEME = {
    "2_orange_green": [RGBColor(255, 192, 0), RGBColor(0, 176, 80)],
    "2_green_orange": [RGBColor(0, 176, 80), RGBColor(255, 192, 0)],
    "2_orange_olive": [RGBColor(255, 192, 0), RGBColor(0, 176, 180)],
    "2_olive_green": [RGBColor(0, 176, 180), RGBColor(0, 176, 80)],
    "2_white_olive": [RGBColor(255, 255, 255), RGBColor(0, 176, 180)],
    "2_blue_lightblue": [RGBColor(0, 176, 240), RGBColor(169, 227, 255)],
    "2_olive_orange": [RGBColor(0, 176, 180), RGBColor(255, 192, 0)],
    "2_gray_lightblue2": [RGBColor(182, 191, 197), RGBColor(169, 227, 255)],
    "1_olive": [RGBColor(0, 176, 180)],
    "1_orange": [RGBColor(255, 192, 0)],
    "1_gray": [RGBColor(182, 191, 197)],
    "1_lightblue": [RGBColor(0, 176, 240)],
    "1_tangerine": [RGBColor(255, 102, 153)],
    "2_tangerine_gray": [RGBColor(255, 102, 153), RGBColor(182, 191, 197)],
    "2_gray_tangerine": [RGBColor(182, 191, 197), RGBColor(255, 102, 153)],
    "2_gray_lightblue": [RGBColor(182, 191, 197), RGBColor(0, 176, 240)],
    "2_green_gray": [RGBColor(0, 176, 80), RGBColor(182, 191, 197)],
    "2_olive_gray": [RGBColor(0, 176, 180), RGBColor(182, 191, 197)],
}


class PPTOperations:
    """
    PPT操作类

    Usage:
        with PPTOperations("output.pptx") as ppt:
            ppt.create_blank_slides(5)
            ppt.add_title(1, content, layout)
            # 退出时自动保存
    """

    def __init__(self, file_path: str | Path, template_path: str | Path | None = None):
        """
        Args:
            file_path: PPT 保存的目标路径
            template_path: 模板路径，如果存在则基于模板加载，否则新建空白 PPT
        """
        self.file_path = str(file_path)
        self.presentation = None

        # 初始化加载逻辑
        if template_path and os.path.exists(template_path):
            logger.info(f"加载模板 PPT: {template_path}")
            self.presentation = Presentation(template_path)
        elif os.path.exists(self.file_path):
            logger.info(f"加载现有 PPT: {self.file_path}")
            self.presentation = Presentation(self.file_path)
        else:
            logger.info("创建新的空白 Presentation 对象")
            self.presentation = Presentation()

    def __enter__(self) -> Self:
        return self

    def __exit__(self, exc_type, exc_val, exc_tb):
        if exc_type is None:
            self.save()
        else:
            logger.error(f"操作过程中发生异常，未保存: {exc_val}")

    def save(self, output_path: str | None = None) -> None:
        """保存 PPT 文件"""
        target_path = output_path or self.file_path
        os.makedirs(os.path.dirname(os.path.abspath(target_path)), exist_ok=True)
        try:
            self.presentation.save(target_path)
            logger.info(f"成功保存 PPT 至: {target_path}")
        except Exception as e:
            logger.error(f"保存 PPT 失败: {e}")
            raise

    def _get_slide(self, page_num: int) -> Slide:
        """
        获取指定页码的幻灯片

        Args:
            page_num (int): 页码，从 1 开始
        """
        total = len(self.presentation.slides)
        if not (1 <= page_num <= total):
            raise IndexError(f"页码 {page_num} 超出范围 (当前总页数: {total})")
        return self.presentation.slides[page_num - 1]

    def init_slides(
        self,
        count: int,
        width_cm: float = 33.867,  # 16:9 宽屏默认值
        height_cm: float = 19.05,
    ) -> None:
        """
        初始化/重置幻灯片尺寸并创建空白页

        Args:
            count (int): 目标总页数
            slide_width (float): 幻灯片宽度（厘米）
            slide_height (float): 幻灯片高度（厘米）
        """
        self.presentation.slide_width = Cm(width_cm)
        self.presentation.slide_height = Cm(height_cm)

        # 确保有足够的页面
        current_count = len(self.presentation.slides)
        needed = count - current_count

        blank_layout = self.presentation.slide_layouts[6]  # 空白布局

        for _ in range(needed):
            self.presentation.slides.add_slide(blank_layout)

        logger.info(f"PPT 已调整为 {len(self.presentation.slides)} 页")

    def add_title(
        self,
        page_num: int,
        content: TextContentModel,
        layout: LayoutModel,
    ) -> None:
        """
        添加标题文本框

        Args:
            page_num (int): 幻灯片页码（从1开始）
            content (TextContentModel): 文本内容配置
            layout (LayoutModel): 布局配置
        """
        slide = self._get_slide(page_num)

        # 1. 使用 layout 对象的数据
        text_box = slide.shapes.add_textbox(
            Cm(layout.left), Cm(layout.top), Cm(layout.width), Cm(layout.height)
        )
        text_frame = text_box.text_frame
        text_frame.text = content.text  # 直接访问属性

        text_frame.word_wrap = content.word_wrap

        # 2. 设置样式
        for paragraph in text_frame.paragraphs:
            paragraph.alignment = layout.alignment.pptx_val

            for run in paragraph.runs:
                run.font.size = Pt(content.font_size)
                run.font.bold = content.font_bold
                run.font.color.rgb = content.font_color.rgb  # 使用 Enum 属性
                pptx_ea_font.set_font(run, content.font_name)

        logger.info(f"Page {page_num}: 添加标题 '{content.text[:10]}...'")

    def add_content(
        self,
        page_num: int,
        content: TextContentModel,
        layout: LayoutModel,
    ) -> None:
        """
        添加内容文本框

        Args:
            page_num (int): 幻灯片页码（从1开始）
            content (TextContentModel): 文本内容配置
            layout (LayoutModel): 布局配置
        """
        slide = self._get_slide(page_num)

        # 1. 使用 layout 对象的数据
        text_box = slide.shapes.add_textbox(
            Cm(layout.left), Cm(layout.top), Cm(layout.width), Cm(layout.height)
        )
        text_frame = text_box.text_frame
        text_frame.text = content.text  # 直接访问属性

        text_frame.word_wrap = content.word_wrap

        # 2. 设置样式
        for paragraph in text_frame.paragraphs:
            paragraph.alignment = layout.alignment.pptx_val

            for run in paragraph.runs:
                run.font.size = Pt(content.font_size)
                run.font.bold = content.font_bold
                run.font.color.rgb = content.font_color.rgb  # 使用 Enum 属性
                pptx_ea_font.set_font(run, content.font_name)

        logger.info(f"Page {page_num}: 添加内容 '{content.text[:10]}...'")

    def add_table(
        self,
        page_num: int,
        layout: LayoutModel,
        data: pd.DataFrame,
        font_name: str = "方正兰亭黑_GBK",
        font_size: int = 10,
    ) -> None:
        """
        添加表格

        Args:
            page_num (int): 幻灯片页码（从1开始）
            layout (LayoutModel): 布局配置
            data (pd.DataFrame): 表格数据
            font_name (str): 字体名称，默认'方正兰亭黑_GBK'
            font_size (int): 字体大小，默认6
        """
        slide = self._get_slide(page_num)
        rows, cols = data.shape

        table_shape = slide.shapes.add_table(
            rows + 1,
            cols,
            Cm(layout.left),
            Cm(layout.top),
            Cm(layout.width),
            Cm(layout.height),
        )
        table = table_shape.table

        # 内部辅助函数：设置单元格样式
        def _set_cell_text(cell, text, is_bold=False):
            cell.text = str(text)
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(font_size)
                    run.font.bold = is_bold
                    try:
                        pptx_ea_font.set_font(run, font_name)
                    except Exception:
                        logger.warning(f"字体 {font_name} 不存在，使用默认字体")

        # 设置表头
        for col_idx, col_name in enumerate(data.columns):
            _set_cell_text(table.cell(0, col_idx), col_name, is_bold=True)

        # 设置数据
        for row_idx in range(rows):
            for col_idx in range(cols):
                val = data.iloc[row_idx, col_idx]
                _set_cell_text(table.cell(row_idx + 1, col_idx), val, is_bold=False)

        logger.info(f"Page {page_num}: 添加 {rows}x{cols} 表格")

    def add_rectangle(
        self,
        page_num: int,
        layout: LayoutModel,
        style: RectangleStyleModel | None = None,
    ) -> None:
        """添加矩形色块

        Args:
            page_num (int): 幻灯片页码（从1开始）
            layout (LayoutModel): 布局配置
            style (RectangleStyleModel): 矩形样式配置
        """
        slide = self._get_slide(page_num)

        # 添加形状
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Cm(layout.left),
            Cm(layout.top),
            Cm(layout.width),
            Cm(layout.height),
        )
        shape.rotation = style.rotation

        # 填充
        shape.fill.solid()
        shape.fill.fore_color.rgb = style.fore_color.rgb  # 直接使用 Enum 属性

        # 边框
        shape.line.color.rgb = style.line_color.rgb
        shape.line.width = Pt(style.line_width)

        if style.is_background:
            shape.fill.background()

        logger.info(f"Page {page_num}: 添加矩形")

    def add_picture(self, page_num: int, img_path: str, layout: LayoutModel) -> None:
        """
        添加图片

        Args:
            page_num (int): 幻灯片页码（从1开始）
            img_path (str): 图片路径
            layout (LayoutModel): 布局配置
        """
        slide = self._get_slide(page_num)
        if not os.path.exists(img_path):
            logger.error(f"图片不存在: {img_path}")
            return

        slide.shapes.add_picture(
            img_path,
            Cm(layout.left),
            Cm(layout.top),
            width=Cm(layout.width),  # 高度可选，按比例缩放
        )
        logger.info(f"Page {page_num}: 添加图片 {os.path.basename(img_path)}")

    def add_chart(
        self,
        page_num: int,
        df: pd.DataFrame,
        layout: LayoutModel,
        config: ChartConfigModel,
    ) -> None:
        """
        添加柱状图

        Args:
            df (pd.DataFrame):
                - Index (索引) -> 变为图例 (系列)
                - Columns (列名) -> 变为 X 轴标签 (类别)
            layout (LayoutModel): Pydantic 布局模型
            config (ChartConfigModel): Pydantic 图表配置模型
        """
        slide = self._get_slide(page_num)

        # 1. 准备数据
        chart_data = ChartData()
        chart_data.categories = df.columns  # 列名作为 X 轴
        for i in range(len(df)):
            # 处理可能存在的 NaN
            series_data = df.iloc[i].fillna(0)
            chart_data.add_series(str(df.index[i]), series_data)

        # 2. 添加图表 (使用 Pydantic layout 对象属性)
        x, y, w, h = (
            Cm(layout.left),
            Cm(layout.top),
            Cm(layout.width),
            Cm(layout.height),
        )

        graphic_frame = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, w, h, chart_data
        )
        chart = graphic_frame.chart

        # 3. 设置颜色
        if config.style_name in CHART_COLORS_THEME:
            colors = CHART_COLORS_THEME[config.style_name]
            for idx, series in enumerate(chart.series):
                if idx < len(colors):
                    series.format.fill.solid()
                    series.format.fill.fore_color.rgb = colors[idx]

        # 字体统一设置 helper
        def _set_axis_font(axis):
            axis.tick_labels.font.size = Pt(config.font_size)
            # axis.tick_labels.font.name = "Arial" # 图表字体通常用英文字体

        # 4. 基础样式设置
        chart.has_title = False
        _set_axis_font(chart.category_axis)

        val_axis = chart.value_axis
        val_axis.visible = config.y_axis_visible
        if config.value_axis_max is not None:
            val_axis.maximum_scale = config.value_axis_max
        _set_axis_font(val_axis)

        # 数据标签
        if config.has_data_labels:
            plot = chart.plots[0]
            plot.has_data_labels = True
            data_labels = plot.data_labels
            data_labels.font.size = Pt(config.font_size)
            data_labels.position = XL_DATA_LABEL_POSITION.CENTER

        # 图例
        chart.has_legend = config.has_legend
        if config.has_legend:
            chart.legend.position = XL_LEGEND_POSITION.TOP
            chart.legend.include_in_layout = False
            chart.legend.font.size = Pt(config.font_size)

        logger.info(f"Page {page_num}: 添加图表")
